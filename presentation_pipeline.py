"""
Presentation Processing Pipeline
Converts PPTX and PDF presentations to structured markdown using native text extraction and AI-powered vision analysis.

Supports corporate presentations, scientific slide decks, and meeting materials.
Uses python-pptx for native PPTX extraction and PyMuPDF/Vision AI for PDFs.
PowerPoint (when installed) enables Vision AI analysis for PPTX files.
"""

import os
import re
import json
import base64
import argparse
import shutil
import tempfile
from pathlib import Path
from typing import List, Dict, Tuple, Optional, Any
from io import BytesIO
import logging
from datetime import datetime

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Load credentials from Windows User environment (with URL validation)
from pipeline_security import load_credentials_from_registry, validate_path, validate_output_path, check_pptx_safe, sanitize_filename
load_credentials_from_registry()


class PresentationPipeline:
    """Pipeline for processing PPTX/PDF presentations into structured markdown."""

    VISION_MODEL = "claude-sonnet-4-6"
    SLIDES_PER_BATCH = 5
    MAX_CONCURRENT_BATCHES = 5
    RENDER_DPI = 150
    CHEMICAL_STRUCTURE_DPI = 250
    MAX_IMAGE_DIMENSION = 1568
    MAX_TOKENS_SLIDE_EXTRACTION = 8192
    MAX_TOKENS_SUMMARY = 4096
    MAX_TOKENS_ANALYSIS = 4096

    MAX_FILE_SIZE_MB = 500  # Reject files larger than 500MB
    CLASSIFICATION_RANKS = {'public': 0, 'confidential': 1, 'secret': 2}

    SMILES_VALID_CHARS = re.compile(r'^[A-Za-z0-9@+\-\[\]()=#$/\\.%:]+$')

    CLASSIFICATION_FILENAME_PATTERNS = {
        'public': [
            r'\bnon[\s\-_]?con\b',
            r'\bpublic\b',
        ],
        'secret': [
            r'\bsecret\b',
            r'\btop[\s\-]?secret\b',
            r'\bstrictly[\s\-]?confidential\b',
        ],
        'confidential': [
            r'\bconfidential\b',
            r'\brestricted\b',
            r'\binternal\s+use\b',
        ],
    }

    CLASSIFICATION_CONTENT_PATTERNS = {
        'public': [
            r'\bfor\s+public\s+release\b',
            r'\bnon[\s\-]?confidential\b',
        ],
        'secret': [
            r'\btop\s+secret\b',
            r'\bstrictly\s+confidential\b',
        ],
        'confidential': [
            r'\bconfidential\b',
            r'\bfor\s+internal\s+use\s+only\b',
            r'\bnot\s+for\s+distribution\b',
            r'\bmerck\s+internal\b',
            r'\bproprietary\b',
            r'\bdo\s+not\s+distribute\b',
            r'\bdo\s+not\s+copy\b',
        ],
    }

    DATE_PATTERNS = [
        (r'^(\d{4})(\d{2})(\d{2})\b', 'prefix_YYYYMMDD'),
        (r'^(\d{4})_(\d{2})_(\d{2})_', 'prefix_YYYY_MM_DD'),
        (r'_(\d{2})\.(\d{2})\.(\d{4})', 'suffix_DD_MM_YYYY'),
        (r'(\d{2})(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)(\d{4})', 'suffix_DDMonYYYY'),
        (r'(January|February|March|April|May|June|July|August|September|October|November|December)\s+(\d{4})', 'text_Month_YYYY'),
        (r'(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\s+(\d{4})', 'text_Mon_YYYY'),
    ]

    BOILERPLATE_PATTERNS = [
        r'(?i)^click\s+to\s+edit',
        r'(?i)^insert\s+(title|subtitle|text)',
        r'(?i)^thank\s+you\s*[!.]?\s*$',
        r'(?i)^questions\s*\??\s*$',
        r'(?i)^(any\s+)?questions\s*(\?|and\s+discussion)',
        r'(?i)^appendix\s*$',
        r'(?i)^backup\s+slides?\s*$',
    ]

    # Pre-compiled patterns for content type classification (single-pass)
    CONTENT_TYPE_PATTERNS = {
        'external': re.compile(r'\b(?:partner|collaboration|alliance|joint|CRO|vendor|external|contractor|outsourc|co-develop)\b', re.IGNORECASE),
        'academic': re.compile(r'\b(?:patient|study|trial|efficacy|endpoint|cohort|publication|abstract|poster|congress|ASCO|AACR|ESMO|hypothesis|method|results|conclusion|p[\s\-]?value)\b', re.IGNORECASE),
        'project': re.compile(r'\b(?:compound|molecule|target|inhibitor|mechanism|discovery|preclinical|clinical|pipeline|candidate|assay|screen|SAR|selectivity|potency|IC50)\b', re.IGNORECASE),
        'operational': re.compile(r'\b(?:update|status|progress|timeline|milestone|deliverable|team|resource|budget|plan|roadmap|strategy|meeting|review|decision|action|department)\b', re.IGNORECASE),
    }

    MONTH_MAP = {
        'jan': 1, 'feb': 2, 'mar': 3, 'apr': 4, 'may': 5, 'jun': 6,
        'jul': 7, 'aug': 8, 'sep': 9, 'oct': 10, 'nov': 11, 'dec': 12,
        'january': 1, 'february': 2, 'march': 3, 'april': 4, 'may': 5, 'june': 6,
        'july': 7, 'august': 8, 'september': 9, 'october': 10, 'november': 11, 'december': 12,
    }

    NAMING_SCHEMES = {
        'default': 'presentation_{title_slug}',
        'dated': '{date}_{classification}_{title_slug}',
        'classified': '{classification}_{date}_{title_slug}',
    }

    def __init__(self, input_folder: str, output_dir: str = "output_presentations",
                 no_vision: bool = False, recursive: bool = False,
                 naming: str = 'default'):
        self.input_folder = Path(input_folder)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.no_vision = no_vision
        self.recursive = recursive
        self.naming = naming

        self.processing_log_path = self.output_dir / "processing_log.txt"
        self._init_processing_log()

        self._existing_files = set(f.stem for f in self.output_dir.glob("*.md"))
        self._client = None

    def _init_processing_log(self):
        if not self.processing_log_path.exists() or self.processing_log_path.stat().st_size == 0:
            with open(self.processing_log_path, 'w', encoding='utf-8') as f:
                f.write("TIMESTAMP\tFILENAME\tFORMAT\tSLIDES\tCLASSIFICATION\tCONTENT_TYPE\t"
                        "CHARS\tQUALITY\tASSESSMENT\tSTATUS\n")

    @staticmethod
    def _tsv_escape(value: str) -> str:
        """Escape a value for safe inclusion in a TSV field."""
        return str(value).replace('\t', ' ').replace('\n', ' ').replace('\r', '')

    def log_processing(self, filename: str, source_format: str, num_slides: int,
                       classification: str, content_type: str, total_chars: int,
                       quality_scores: Optional[Dict], status: str):
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        quality = quality_scores['overall'] if quality_scores else ''
        assessment = self._tsv_escape(quality_scores['assessment']) if quality_scores else ''
        entry = (f"{timestamp}\t{self._tsv_escape(filename)}\t{source_format}\t{num_slides}\t"
                 f"{classification}\t{content_type}\t{total_chars}\t{quality}\t{assessment}\t{status}\n")
        with open(self.processing_log_path, 'a', encoding='utf-8') as f:
            f.write(entry)

    API_MAX_RETRIES = 3
    API_RETRY_DELAY = 2.0  # seconds, doubled each retry

    def get_anthropic_client(self):
        if self._client is not None:
            return self._client

        from anthropic import Anthropic
        auth_token = os.environ.get('ANTHROPIC_AUTH_TOKEN')
        base_url = os.environ.get('ANTHROPIC_BASE_URL')

        if auth_token and base_url:
            self._client = Anthropic(api_key=auth_token, base_url=base_url)
        elif os.environ.get('ANTHROPIC_API_KEY'):
            self._client = Anthropic(api_key=os.environ['ANTHROPIC_API_KEY'])
        else:
            logger.error("API credentials not found")
            return None

        return self._client

    def _api_call_with_retry(self, call_fn, description: str = "API call"):
        """Execute an API call with exponential backoff retry on transient errors."""
        import time
        delay = self.API_RETRY_DELAY
        for attempt in range(self.API_MAX_RETRIES):
            try:
                return call_fn()
            except Exception as e:
                error_str = str(e)
                is_transient = any(s in error_str for s in ('429', '503', '529', 'timeout', 'Timeout', 'overloaded'))
                if is_transient and attempt < self.API_MAX_RETRIES - 1:
                    logger.warning(f"  {description} attempt {attempt + 1} failed ({error_str}), retrying in {delay:.0f}s...")
                    time.sleep(delay)
                    delay *= 2
                else:
                    raise
        return None  # unreachable

    def list_presentation_files(self) -> List[Path]:
        files = []
        if self.recursive:
            files.extend(self.input_folder.rglob("*.pptx"))
            files.extend(self.input_folder.rglob("*.pdf"))
        else:
            files.extend(self.input_folder.glob("*.pptx"))
            files.extend(self.input_folder.glob("*.pdf"))
        files = sorted(files, key=lambda p: p.name.lower())
        logger.info(f"Found {len(files)} presentation files ({sum(1 for f in files if f.suffix == '.pptx')} PPTX, "
                    f"{sum(1 for f in files if f.suffix == '.pdf')} PDF)")
        return files

    def parse_filename(self, filepath: Path) -> Dict:
        stem = filepath.stem
        result = {
            'filename': filepath.name,
            'source_format': filepath.suffix.lstrip('.').lower(),
            'date': None,
            'title': stem,
            'no_structures': bool(re.search(r'\[No\s+Structures\]', stem, re.IGNORECASE)),
        }

        for pattern, fmt in self.DATE_PATTERNS:
            match = re.search(pattern, stem)
            if match:
                try:
                    if fmt == 'prefix_YYYYMMDD':
                        result['date'] = f"{match.group(1)}-{match.group(2)}-{match.group(3)}"
                        result['title'] = stem[match.end():].strip(' _-')
                    elif fmt == 'prefix_YYYY_MM_DD':
                        result['date'] = f"{match.group(1)}-{match.group(2)}-{match.group(3)}"
                        result['title'] = stem[match.end():].strip(' _-')
                    elif fmt == 'suffix_DD_MM_YYYY':
                        result['date'] = f"{match.group(3)}-{match.group(2)}-{match.group(1)}"
                        title_end = match.start()
                        result['title'] = stem[:title_end].strip(' _-')
                    elif fmt == 'suffix_DDMonYYYY':
                        day = int(match.group(1))
                        month = self.MONTH_MAP[match.group(2).lower()]
                        year = int(match.group(3))
                        result['date'] = f"{year:04d}-{month:02d}-{day:02d}"
                        title_end = match.start()
                        result['title'] = stem[:title_end].strip(' _-')
                    elif fmt in ('text_Month_YYYY', 'text_Mon_YYYY'):
                        month = self.MONTH_MAP[match.group(1).lower()]
                        year = int(match.group(2))
                        result['date'] = f"{year:04d}-{month:02d}-01"
                except (ValueError, KeyError):
                    pass
                break

        result['title'] = re.sub(r'^(Confidential\s*(\[.*?\])?\s*_?)', '', result['title'], flags=re.IGNORECASE).strip(' _-')
        result['title'] = re.sub(r'_+', ' ', result['title']).strip()
        if not result['title']:
            result['title'] = stem

        return result

    def extract_pptx_content(self, pptx_path: Path) -> Dict:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Safety check: reject PPTX with macros/ActiveX
        if not check_pptx_safe(pptx_path):
            logger.error(f"PPTX safety check failed (macros/ActiveX detected): {pptx_path.name}")
            return {'slides': [], 'metadata': {}, 'num_slides': 0}

        prs = Presentation(str(pptx_path))

        metadata = {}
        props = prs.core_properties
        if props.title:
            metadata['title'] = props.title
        if props.author:
            metadata['author'] = props.author
        if props.last_modified_by:
            metadata['last_modified_by'] = props.last_modified_by
        if props.created:
            metadata['created'] = props.created.strftime('%Y-%m-%d') if props.created else None
        if props.modified:
            metadata['modified'] = props.modified.strftime('%Y-%m-%d') if props.modified else None
        if props.subject:
            metadata['subject'] = props.subject
        if props.revision:
            metadata['revision'] = props.revision

        slides_data = []
        all_hyperlinks = []
        has_speaker_notes = False
        media_flags = []
        total_images = 0
        total_charts = 0

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_data = {
                'number': slide_idx,
                'title': None,
                'text_blocks': [],
                'tables': [],
                'notes': None,
                'is_boilerplate': False,
                'has_images': False,
                'has_charts': False,
            }

            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_data['title'] = slide.shapes.title.text.strip()

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    if table_data:
                        slide_data['tables'].append(table_data)

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_data['has_images'] = True
                    total_images += 1

                elif shape.shape_type in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.IGX_GRAPHIC, MSO_SHAPE_TYPE.DIAGRAM):
                    slide_data['has_charts'] = True
                    total_charts += 1

                elif shape.shape_type in (MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT):
                    media_flags.append({
                        'slide': slide_idx,
                        'type': str(shape.shape_type),
                        'name': shape.name if hasattr(shape, 'name') else 'unknown',
                    })

                elif shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_data['text_blocks'].append(para_text)
                        for run in paragraph.runs:
                            if run.hyperlink and run.hyperlink.address:
                                all_hyperlinks.append({
                                    'slide': slide_idx,
                                    'text': run.text.strip(),
                                    'url': run.hyperlink.address,
                                })

            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text and not re.match(r'^(\d+|slide\s*\d*)$', notes_text, re.IGNORECASE):
                        slide_data['notes'] = notes_text
                        has_speaker_notes = True
            except Exception:
                pass

            all_text = ' '.join(slide_data['text_blocks'])
            if slide_data['title']:
                all_text = slide_data['title'] + ' ' + all_text
            for pattern in self.BOILERPLATE_PATTERNS:
                if re.search(pattern, all_text):
                    slide_data['is_boilerplate'] = True
                    break

            slides_data.append(slide_data)

        has_visual_content = total_images > 0 or total_charts > 0
        if has_visual_content:
            logger.info(f"  Visual content detected: {total_images} images, {total_charts} charts/diagrams")

        return {
            'metadata': metadata,
            'slides': slides_data,
            'hyperlinks': all_hyperlinks,
            'media_flags': media_flags,
            'has_speaker_notes': has_speaker_notes,
            'num_slides': len(slides_data),
            'extraction_method': 'native_text',
            'has_visual_content': has_visual_content,
            'image_count': total_images,
            'chart_count': total_charts,
        }

    def extract_pdf_content(self, pdf_path: Path) -> Dict:
        import fitz
        import pdfplumber

        slides_data = []
        total_chars = 0
        num_pages = 0
        garbled_pages = []

        try:
            doc = fitz.open(str(pdf_path))
            num_pages = len(doc)

            for page_idx in range(num_pages):
                page = doc[page_idx]
                blocks = page.get_text('dict')['blocks']

                text_lines = []
                has_rotated = False
                for b in blocks:
                    if 'lines' not in b:
                        continue
                    for line in b['lines']:
                        d = line['dir']
                        line_text = ''.join(s['text'] for s in line['spans']).strip()
                        if not line_text:
                            continue
                        is_rotated = abs(d[0]) < 0.5
                        if is_rotated:
                            has_rotated = True
                            text_lines.append(f"[Axis label: {line_text}]")
                        else:
                            text_lines.append(line_text)

                text_lines = self._filter_chart_axis_data(text_lines)
                page_text = '\n'.join(text_lines)
                total_chars += len(page_text)

                slide_data = {
                    'number': page_idx + 1,
                    'title': text_lines[0][:100] if text_lines else None,
                    'text_blocks': text_lines,
                    'tables': [],
                    'notes': None,
                    'is_boilerplate': False,
                }

                if has_rotated or self._has_garbled_text(text_lines):
                    garbled_pages.append(page_idx)

                slides_data.append(slide_data)

            doc.close()

            # Use pdfplumber for table extraction only (superior table detection)
            with pdfplumber.open(str(pdf_path)) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    if tables and page_idx < len(slides_data):
                        for table in tables:
                            cleaned = [[cell or '' for cell in row] for row in table if row]
                            if cleaned:
                                slides_data[page_idx]['tables'].append(cleaned)
                        # Deduplicate: remove text_blocks that are subsets of table content
                        slides_data[page_idx]['text_blocks'] = self._deduplicate_text_vs_tables(
                            slides_data[page_idx]['text_blocks'],
                            slides_data[page_idx]['tables']
                        )

        except Exception as e:
            logger.error(f"PDF text extraction failed: {e}")
            return {'slides': [], 'num_slides': 0, 'extraction_method': 'failed'}

        avg_chars = total_chars / max(num_pages, 1)
        needs_vision_sparse = avg_chars < 50
        needs_vision_garbled = len(garbled_pages) > num_pages * 0.2

        extraction_method = 'native_text'
        if (needs_vision_sparse or needs_vision_garbled) and not self.no_vision:
            reason = f"sparse text ({avg_chars:.0f} chars/page)" if needs_vision_sparse else f"garbled content on {len(garbled_pages)} pages"
            logger.info(f"  PDF quality issue: {reason} — using Vision AI enhancement")
            extraction_method = 'native_vision_enhanced'
            vision_content = self._extract_pdf_with_vision(pdf_path, num_pages)
            if vision_content:
                slides_data = vision_content

        return {
            'metadata': {},
            'slides': slides_data,
            'hyperlinks': [],
            'media_flags': [],
            'has_speaker_notes': False,
            'num_slides': num_pages,
            'extraction_method': extraction_method,
        }

    def _extract_pdf_with_vision(self, pdf_path: Path, num_pages: int) -> Optional[List[Dict]]:
        slide_images = self._convert_pdf_to_images(pdf_path)
        if not slide_images:
            return None

        encoded_images = []
        for img in slide_images:
            b64, media_type = self.encode_image_to_base64(img)
            if b64:
                encoded_images.append((b64, media_type))

        if not encoded_images:
            return None

        all_slides = []
        batches = [encoded_images[i:i+self.SLIDES_PER_BATCH]
                   for i in range(0, len(encoded_images), self.SLIDES_PER_BATCH)]

        for batch_idx, batch in enumerate(batches):
            batch_start = batch_idx * self.SLIDES_PER_BATCH
            result = self._extract_slides_batch_vision(batch, batch_start)
            all_slides.extend(result)

        return all_slides

    def _convert_pdf_to_images(self, pdf_path: Path) -> List[Any]:
        try:
            import fitz
            from PIL import Image

            Image.MAX_IMAGE_PIXELS = 300_000_000  # ~300MP, generous but bounded
            doc = fitz.open(str(pdf_path))
            images = []
            zoom = self.RENDER_DPI / 72
            mat = fitz.Matrix(zoom, zoom)

            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(matrix=mat)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                images.append(img)

            doc.close()
            return images
        except Exception as e:
            logger.error(f"Error converting PDF to images: {e}")
            return []

    def _filter_chart_axis_data(self, lines: List[str]) -> List[str]:
        """Remove sequences of pure-numeric lines that are chart axis tick marks."""
        filtered = []
        i = 0
        while i < len(lines):
            # Check for a run of 3+ purely numeric lines (chart axis values)
            if re.match(r'^[\d.,\s\-]+$', lines[i]) and len(lines[i].strip()) < 20:
                run_start = i
                while i < len(lines) and re.match(r'^[\d.,\s\-]+$', lines[i]) and len(lines[i].strip()) < 20:
                    i += 1
                if i - run_start >= 3:
                    # Collapse to a single marker
                    values = [lines[j].strip() for j in range(run_start, i)]
                    filtered.append(f"[Chart axis: {', '.join(values)}]")
                else:
                    # Not enough to be an axis — keep the lines
                    for j in range(run_start, i):
                        filtered.append(lines[j])
            else:
                filtered.append(lines[i])
                i += 1
        return filtered

    def _has_garbled_text(self, lines: List[str]) -> bool:
        """Detect signs of garbled extraction: reversed words, fragmented subscripts."""
        garbled_indicators = 0
        for line in lines:
            # Reversed common scientific terms
            if re.search(r'\b(ytisnetni|gniniameR|nietorP|evitaleR|noitartnecnoC)\b', line):
                garbled_indicators += 2
            # Subscript fragments: line is just "50", "max", "1/2", "dss" after a proper term
            if re.match(r'^(50|max|1/2|dss|0-last|min|ss)\s*$', line):
                garbled_indicators += 1
        return garbled_indicators >= 3

    def _deduplicate_text_vs_tables(self, text_blocks: List[str], tables: List[List[List[str]]]) -> List[str]:
        """Remove text blocks that are duplicated in table content."""
        if not tables:
            return text_blocks

        # Build set of all cell content from tables
        table_content = set()
        for table in tables:
            for row in table:
                for cell in row:
                    if cell:
                        # Add both the full cell and individual lines within multi-line cells
                        table_content.add(cell.strip())
                        for line in cell.split('\n'):
                            if line.strip():
                                table_content.add(line.strip())

        # Pre-filter: only check substring against cells longer than the block
        # Sort table_content by length for efficient substring checking
        long_cells = sorted([tc for tc in table_content if len(tc) > 10], key=len, reverse=True)

        # Filter text blocks, keeping those not found in tables
        filtered = []
        for block in text_blocks:
            block_stripped = block.strip()
            if not block_stripped:
                continue
            if block_stripped in table_content:
                continue
            # Only check substring against cells that are longer than the block
            is_in_table = any(block_stripped in tc for tc in long_cells if len(tc) > len(block_stripped))
            if not is_in_table:
                filtered.append(block)

        return filtered

    def _detect_recurring_footer(self, slides_data: List[Dict]) -> Optional[List[str]]:
        """Detect text that appears on 3+ slides — likely a footer/header bar."""
        from collections import Counter
        block_counts = Counter()
        for slide in slides_data:
            seen = set()
            for block in slide.get('text_blocks', []):
                b = block.strip()
                if b and b not in seen and len(b) > 5:
                    seen.add(b)
                    block_counts[b] += 1
        threshold = max(3, len(slides_data) * 0.15)
        footers = [text for text, count in block_counts.items() if count >= threshold]
        return footers if footers else None

    def _merge_short_fragments(self, text_blocks: List[str]) -> List[str]:
        """Merge consecutive short text fragments that form logical units (e.g. stat blocks)."""
        if not text_blocks:
            return text_blocks
        merged = []
        buffer = []
        for block in text_blocks:
            stripped = block.strip()
            if len(stripped) <= 35 and not stripped.startswith('[') and not stripped.startswith('|'):
                buffer.append(stripped)
            else:
                if buffer:
                    merged.append(' '.join(buffer))
                    buffer = []
                merged.append(block)
        if buffer:
            merged.append(' '.join(buffer))
        return merged

    def encode_image_to_base64(self, image, format="JPEG") -> Tuple[Optional[str], Optional[str]]:
        try:
            from PIL import Image

            width, height = image.size
            if max(width, height) > self.MAX_IMAGE_DIMENSION:
                scale = self.MAX_IMAGE_DIMENSION / max(width, height)
                image = image.resize((int(width * scale), int(height * scale)), Image.Resampling.LANCZOS)

            if image.mode == 'RGBA':
                image = image.convert('RGB')

            buffered = BytesIO()
            image.save(buffered, format=format, quality=85)
            img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
            media_type = f"image/{format.lower()}"
            return img_base64, media_type
        except Exception as e:
            logger.error(f"Error encoding image: {e}")
            return None, None

    def _extract_slides_batch_vision(self, encoded_images: List[Tuple[str, str]],
                                      batch_start: int) -> List[Dict]:
        client = self.get_anthropic_client()
        if not client:
            return []

        content_blocks = []
        for i, (img_b64, media_type) in enumerate(encoded_images):
            content_blocks.append({
                "type": "image",
                "source": {"type": "base64", "media_type": media_type, "data": img_b64}
            })
            content_blocks.append({
                "type": "text",
                "text": f"[Slide {batch_start + i + 1}]"
            })

        prompt = """For each slide image above, extract ALL text content and describe visual elements.
Return a JSON array with one object per slide:
[
  {
    "slide_number": 1,
    "title": "slide title or null",
    "text_content": "all text on the slide",
    "visual_elements": "description of charts, diagrams, images",
    "is_boilerplate": false
  }
]
Return ONLY valid JSON."""

        content_blocks.append({"type": "text", "text": prompt})

        try:
            response = self._api_call_with_retry(
                lambda: client.messages.create(
                    model=self.VISION_MODEL,
                    max_tokens=self.MAX_TOKENS_SLIDE_EXTRACTION,
                    messages=[{"role": "user", "content": content_blocks}]
                ),
                description="Vision batch extraction"
            )
            results = self.parse_json_response(response.content[0].text)
            if isinstance(results, list):
                slides = []
                for item in results:
                    text_content = item.get('text_content', '')
                    slides.append({
                        'number': item.get('slide_number', batch_start + len(slides) + 1),
                        'title': item.get('title'),
                        'text_blocks': [line.strip() for line in text_content.split('\n') if line.strip()],
                        'tables': [],
                        'notes': None,
                        'is_boilerplate': item.get('is_boilerplate', False),
                        'visual_elements': item.get('visual_elements'),
                    })
                return slides
        except Exception as e:
            logger.error(f"Vision batch extraction failed: {e}")

        return []

    def detect_classification(self, filename: str, all_text: str) -> Tuple[str, List[str]]:
        signals = []
        detected_level = None

        filename_lower = filename.lower()
        for level in ('public', 'secret', 'confidential'):
            for pattern in self.CLASSIFICATION_FILENAME_PATTERNS[level]:
                if re.search(pattern, filename_lower):
                    signals.append(f"filename: matched '{pattern}' -> {level}")
                    if detected_level is None or self._classification_rank(level) > self._classification_rank(detected_level):
                        detected_level = level
                    break

        text_lower = all_text.lower()[:10000]
        for level in ('secret', 'confidential', 'public'):
            matched = False
            for pattern in self.CLASSIFICATION_CONTENT_PATTERNS[level]:
                match = re.search(pattern, text_lower)
                if match:
                    signals.append(f"content: matched '{match.group(0)}' -> {level}")
                    if detected_level is None or self._classification_rank(level) > self._classification_rank(detected_level):
                        detected_level = level
                    matched = True
                    break
            if matched:
                break

        if detected_level is None:
            detected_level = 'confidential'
            signals.append("default: no explicit markers found, defaulting to confidential")

        return detected_level, signals

    @classmethod
    def _classification_rank(cls, level: str) -> int:
        return cls.CLASSIFICATION_RANKS.get(level, 1)

    def classify_content_type(self, slides_data: List[Dict], filename: str) -> str:
        if re.search(r'\bagenda\b', filename, re.IGNORECASE):
            return 'agenda'

        text_parts = []
        sentences_per_slide = []
        for slide in slides_data:
            slide_text = ' '.join(slide.get('text_blocks', []))
            text_parts.append(slide_text)
            sentence_count = len(re.findall(r'[.!?]\s', slide_text)) + (1 if slide_text.strip() else 0)
            sentences_per_slide.append(sentence_count)

        total_text = ' '.join(text_parts)
        if not total_text.strip():
            return 'template'

        non_boilerplate = [s for s in slides_data if not s.get('is_boilerplate')]
        if len(non_boilerplate) <= 1:
            return 'template'

        if self.no_vision:
            few_sentence_slides = sum(1 for c in sentences_per_slide if c < 3)
            avg_chars = len(total_text) / max(len(non_boilerplate), 1)
            if (len(sentences_per_slide) > 0
                    and few_sentence_slides / len(sentences_per_slide) > 0.7
                    and avg_chars < 200):
                return 'agenda'

        # Single-pass content type classification using pre-compiled patterns
        ext_count = len(self.CONTENT_TYPE_PATTERNS['external'].findall(total_text))
        acad_count = len(self.CONTENT_TYPE_PATTERNS['academic'].findall(total_text))
        proj_count = len(self.CONTENT_TYPE_PATTERNS['project'].findall(total_text))
        op_count = len(self.CONTENT_TYPE_PATTERNS['operational'].findall(total_text))

        scores = {
            'external': ext_count,
            'project': proj_count,
            'academic': acad_count,
            'operational': op_count,
        }
        best = max(scores, key=scores.get)
        best_score = scores[best]

        if best_score < 2:
            return 'operational'

        return best

    @staticmethod
    def _empty_vision_result() -> Dict:
        return {'figures': [], 'chemical_structures': [], 'content_type_hint': None,
                'audience': None, 'topics': []}

    def analyze_with_vision(self, filepath: Path, slides_data: List[Dict],
                            source_format: str, skip: bool = False,
                            preloaded_images: Optional[List[Any]] = None) -> Dict:
        if self.no_vision or skip:
            return self._empty_vision_result()

        client = self.get_anthropic_client()
        if not client:
            return self._empty_vision_result()

        if preloaded_images is not None:
            images = preloaded_images
        elif source_format == 'pptx':
            images = self._convert_pptx_to_images(filepath)
        else:
            images = self._convert_pdf_to_images(filepath)

        if not images:
            return self._empty_vision_result()

        sample_indices = self._select_representative_slides(slides_data, images)
        sample_images = [(images[i], i + 1) for i in sample_indices if i < len(images)]

        content_blocks = []
        for img, slide_num in sample_images:
            b64, media_type = self.encode_image_to_base64(img)
            if b64:
                content_blocks.append({
                    "type": "image",
                    "source": {"type": "base64", "media_type": media_type, "data": b64}
                })
                content_blocks.append({"type": "text", "text": f"[Slide {slide_num}]"})

        if not content_blocks:
            return self._empty_vision_result()

        prompt = """Analyze these presentation slides. Return a JSON object with:
{
  "content_type": "external" or "project" or "academic" or "operational" or "agenda" or "template",
  "audience": "internal_team" or "leadership" or "external_partner" or "conference",
  "topics": ["topic1", "topic2"],
  "figures": [
    {
      "slide": 1,
      "label": "Figure 1",
      "type": "bar_chart|line_chart|pie_chart|scatter_plot|flow_diagram|chemical_structure|reaction_scheme|table_figure|image|other",
      "title": "figure title",
      "description": "what the figure shows",
      "key_findings": "main takeaway from this figure",
      "x_axis": "axis label or null",
      "y_axis": "axis label or null",
      "smiles": "SMILES notation for chemical structures or null",
      "smiles_confidence": "high|medium|low or null"
    }
  ],
  "classification_visual_cues": ["any confidentiality banners, watermarks, or stamps observed"]
}

For chemical structures: provide SMILES notation with confidence level.
For figures: describe data, axes, legends, and key findings.
Topics should be domain-specific tags (e.g., "oncology", "immunology", "operations", "discovery").
Return ONLY valid JSON."""

        content_blocks.append({"type": "text", "text": prompt})

        try:
            response = self._api_call_with_retry(
                lambda: client.messages.create(
                    model=self.VISION_MODEL,
                    max_tokens=self.MAX_TOKENS_ANALYSIS,
                    messages=[{"role": "user", "content": content_blocks}]
                ),
                description="Vision analysis"
            )
            result = self.parse_json_response(response.content[0].text)
            if isinstance(result, dict):
                figures = result.get('figures', [])
                chem_structures = [f for f in figures
                                   if f.get('type') in ('chemical_structure', 'reaction_scheme')]
                return {
                    'figures': [f for f in figures
                               if f.get('type') not in ('chemical_structure', 'reaction_scheme')],
                    'chemical_structures': chem_structures,
                    'content_type_hint': result.get('content_type'),
                    'audience': result.get('audience'),
                    'topics': result.get('topics', []),
                    'classification_visual_cues': result.get('classification_visual_cues', []),
                }
        except Exception as e:
            logger.error(f"Vision analysis failed: {e}")

        return self._empty_vision_result()

    def _convert_pptx_to_images(self, pptx_path: Path) -> List[Any]:
        """Convert PPTX slides to images using PowerPoint COM automation."""
        if not self._check_powerpoint_available():
            return []

        from PIL import Image
        import win32com.client
        import pythoncom

        images = []
        tmp_dir = None
        ppt_app = None

        try:
            pythoncom.CoInitialize()
            tmp_dir = tempfile.mkdtemp(prefix="pptx_slides_")
            abs_path = str(pptx_path.resolve())

            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            presentation = ppt_app.Presentations.Open(abs_path, WithWindow=False)

            try:
                for i, slide in enumerate(presentation.Slides):
                    img_path = os.path.join(tmp_dir, f"slide_{i+1:03d}.png")
                    slide.Export(img_path, "PNG", 1920, 1080)
                    if os.path.exists(img_path):
                        img = Image.open(img_path)
                        images.append(img.copy())
                        img.close()
            finally:
                try:
                    presentation.Close()
                except Exception:
                    pass

            logger.info(f"  Exported {len(images)} slide images via PowerPoint")

        except Exception as e:
            logger.error(f"  PowerPoint slide export failed: {e}")
        finally:
            if ppt_app:
                try:
                    ppt_app.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()
            # Clean up temp files
            if tmp_dir:
                shutil.rmtree(tmp_dir, ignore_errors=True)

        return images

    def _check_powerpoint_available(self) -> bool:
        """Check if PowerPoint is installed and accessible via COM."""
        if hasattr(self, '_powerpoint_available'):
            return self._powerpoint_available

        try:
            import win32com.client
            import pythoncom
            pythoncom.CoInitialize()
            ppt = None
            try:
                ppt = win32com.client.Dispatch("PowerPoint.Application")
                self._powerpoint_available = True
            except Exception as e:
                logger.warning(f"  PowerPoint not available for slide rendering: {e}")
                self._powerpoint_available = False
            finally:
                if ppt:
                    try:
                        ppt.Quit()
                    except Exception:
                        pass
                pythoncom.CoUninitialize()
            return self._powerpoint_available
        except ImportError:
            logger.warning("  win32com not installed — PowerPoint rendering unavailable. "
                          "Install with: pip install pywin32")
            self._powerpoint_available = False
            return False

    def _select_representative_slides(self, slides_data: List[Dict],
                                       images: List[Any]) -> List[int]:
        max_slides = min(15, len(images))
        if len(images) <= max_slides:
            return list(range(len(images)))

        non_boilerplate = [i for i, s in enumerate(slides_data)
                          if not s.get('is_boilerplate') and i < len(images)]

        if len(non_boilerplate) <= max_slides:
            return non_boilerplate

        step = len(non_boilerplate) / max_slides
        return [non_boilerplate[int(i * step)] for i in range(max_slides)]

    def refine_chemical_structures(self, filepath: Path, structures: List[Dict],
                                    slides_data: List[Dict], source_format: str,
                                    preloaded_images: Optional[List[Any]] = None) -> List[Dict]:
        if not structures or self.no_vision:
            return structures

        low_confidence = [s for s in structures
                         if s.get('smiles_confidence') in ('low', 'medium', None)]
        if not low_confidence:
            return structures

        logger.info(f"  Refining SMILES for {len(low_confidence)} structures...")
        client = self.get_anthropic_client()
        if not client:
            return structures

        if preloaded_images is not None:
            images = preloaded_images
        elif source_format == 'pptx':
            images = self._convert_pptx_to_images(filepath)
        else:
            images = self._convert_pdf_to_images(filepath)

        if not images:
            return structures

        refined_count = 0
        for struct in low_confidence:
            slide_num = struct.get('slide', 1) - 1
            if slide_num >= len(images):
                continue

            img = images[slide_num]
            from PIL import Image
            width, height = img.size
            scale = self.CHEMICAL_STRUCTURE_DPI / self.RENDER_DPI
            new_size = (int(width * scale), int(height * scale))
            img_hires = img.resize(new_size, Image.Resampling.LANCZOS)

            b64, media_type = self.encode_image_to_base64(img_hires)
            if not b64:
                continue

            slide_text = ''
            if slide_num < len(slides_data):
                slide_text = ' '.join(slides_data[slide_num].get('text_blocks', []))[:500]

            label = struct.get('label', 'chemical structure')
            prompt = f"""Focus on the chemical structure labeled "{label}" on this slide.

Surrounding text context: {slide_text}

Provide ONLY a JSON object with:
- "smiles": Canonical SMILES notation. Include stereochemistry where visible.
- "smiles_confidence": "high" if simple/clear, "medium" if moderate complexity, "low" if uncertain
- "molecular_formula": e.g. "C23H28BrN7O" if determinable
- "molecular_weight": approximate MW if formula is known, null otherwise

Return ONLY valid JSON."""

            content = [
                {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": b64}},
                {"type": "text", "text": prompt}
            ]

            try:
                response = self._api_call_with_retry(
                    lambda: client.messages.create(
                        model=self.VISION_MODEL,
                        max_tokens=1024,
                        messages=[{"role": "user", "content": content}]
                    ),
                    description=f"SMILES refinement ({label})"
                )
                result = self.parse_json_response(response.content[0].text)
                if result and result.get('smiles') and self.validate_smiles(result['smiles']):
                    struct['smiles'] = result['smiles']
                    struct['smiles_confidence'] = result.get('smiles_confidence', 'medium')
                    if result.get('molecular_formula'):
                        struct['molecular_formula'] = result['molecular_formula']
                    if result.get('molecular_weight'):
                        struct['molecular_weight'] = result['molecular_weight']
                    refined_count += 1
            except Exception as e:
                logger.error(f"  SMILES refinement failed for {label}: {e}")

        logger.info(f"  SMILES refinement: {refined_count}/{len(low_confidence)} improved")
        return structures

    def validate_smiles(self, smiles: str) -> bool:
        if not smiles or len(smiles) < 3:
            return False
        if not self.SMILES_VALID_CHARS.match(smiles):
            return False
        if smiles.count('(') != smiles.count(')'):
            return False
        if smiles.count('[') != smiles.count(']'):
            return False
        return True

    def _detect_language(self, slides_data: List[Dict]) -> str:
        """Detect presentation language from slide text. Returns ISO 639-1 code.

        Currently supports English and German only; defaults to 'en' for other languages.
        """
        text_sample = []
        for slide in slides_data[:20]:
            text_sample.extend(slide.get('text_blocks', []))
        sample = ' '.join(text_sample)[:3000].lower()

        words = re.findall(r'\b[a-zäöüß]+\b', sample)
        if not words:
            return 'en'

        german_indicators = {'und', 'der', 'die', 'das', 'ist', 'für', 'mit', 'auf',
                             'nicht', 'ein', 'eine', 'wird', 'bei', 'nach', 'auch',
                             'sich', 'von', 'oder', 'dem', 'den', 'des', 'wie',
                             'sind', 'werden', 'kann', 'über', 'zur', 'zum', 'alle',
                             'bitte', 'hier', 'neue', 'mehr', 'wenn', 'noch'}
        english_indicators = {'the', 'and', 'for', 'with', 'that', 'this', 'from',
                              'are', 'was', 'have', 'has', 'been', 'will', 'can',
                              'not', 'but', 'they', 'which', 'their', 'would',
                              'should', 'could', 'than', 'other', 'into', 'these'}

        german_count = sum(1 for w in words if w in german_indicators)
        english_count = sum(1 for w in words if w in english_indicators)
        total_indicator_words = german_count + english_count

        if total_indicator_words == 0:
            return 'en'

        german_ratio = german_count / total_indicator_words
        if german_ratio > 0.4:
            return 'de'
        return 'en'

    def generate_summary_and_takeaways(self, slides_data: List[Dict],
                                        file_info: Dict, language: str = 'en') -> Tuple[str, str]:
        if self.no_vision:
            return '', ''

        client = self.get_anthropic_client()
        if not client:
            return '', ''

        slide_text = []
        for slide in slides_data:
            if slide.get('is_boilerplate'):
                continue
            title = slide.get('title') or ''
            text = ' '.join(slide.get('text_blocks', []))
            if title or text:
                slide_text.append(f"Slide {slide['number']}: {title}\n{text}")

        combined = '\n\n'.join(slide_text)[:12000]

        language_instruction = ""
        if language != 'en':
            language_instruction = ("\n\nIMPORTANT: The presentation content is in a non-English language. "
                                   "Write your summary and takeaways entirely in English. "
                                   "Translate key points and findings into clear English.\n")

        prompt = f"""Based on this presentation content, provide:
1. An executive summary (2-4 paragraphs) covering the main points, findings, and conclusions.
2. Key takeaways (4-7 bullet points starting with "- ").
{language_instruction}
Presentation: "{file_info.get('title', 'Untitled')}"

Content:
{combined}

Format your response as:
[SUMMARY]
(executive summary here)

[TAKEAWAYS]
(bullet points here)"""

        try:
            response = self._api_call_with_retry(
                lambda: client.messages.create(
                    model=self.VISION_MODEL,
                    max_tokens=self.MAX_TOKENS_SUMMARY,
                    messages=[{"role": "user", "content": prompt}]
                ),
                description="Summary generation"
            )
            text = response.content[0].text

            summary = ''
            takeaways = ''
            if '[SUMMARY]' in text and '[TAKEAWAYS]' in text:
                parts = text.split('[TAKEAWAYS]')
                summary = parts[0].replace('[SUMMARY]', '').strip()
                takeaways = parts[1].strip() if len(parts) > 1 else ''
            elif '[SUMMARY]' in text:
                summary = text.replace('[SUMMARY]', '').strip()
            else:
                summary = text.strip()

            return summary, takeaways
        except Exception as e:
            logger.error(f"Summary generation failed: {e}")
            return '', ''

    def extract_action_items(self, slides_data: List[Dict]) -> List[Dict]:
        action_items = []
        action_title_patterns = [
            r'(?i)\b(next\s+steps?|action\s+items?|to[\s-]?do)\b',
        ]

        for slide in slides_data:
            title = slide.get('title') or ''

            is_action_slide = any(re.search(p, title) for p in action_title_patterns)
            if not is_action_slide:
                continue

            for block in slide.get('text_blocks', []):
                block = block.strip()
                if len(block) > 10 and not re.match(r'^(next\s+steps?|action\s+items?|to[\s-]?do)', block, re.IGNORECASE):
                    item = {'action': block, 'slide': slide['number'], 'owner': None, 'deadline': None}

                    owner_match = re.search(r'\b([A-Z][a-z]+ [A-Z][a-z]+)\b', block)
                    if owner_match:
                        item['owner'] = owner_match.group(1)

                    date_match = re.search(r'\b(\d{1,2}[./]\d{1,2}[./]\d{2,4}|\d{4}-\d{2}-\d{2}|Q[1-4]\s*\d{4})\b', block)
                    if date_match:
                        item['deadline'] = date_match.group(1)

                    action_items.append(item)

        return action_items

    def extract_key_metrics(self, slides_data: List[Dict]) -> List[Dict]:
        metrics = []
        metric_patterns = [
            (r'(\d+(?:\.\d+)?)\s*%', 'percentage'),
            (r'(?:n\s*=\s*|N\s*=\s*)(\d+)', 'sample_size'),
            (r'(?:p\s*[<>=]\s*)(\d+\.\d+)', 'p_value'),
            (r'(?:HR|OR|RR)\s*[=:]\s*(\d+\.\d+)', 'ratio'),
            (r'(\d+(?:\.\d+)?)\s*(?:mg|µg|ng|mL|µL)', 'dosage'),
        ]

        for slide in slides_data:
            if slide.get('is_boilerplate'):
                continue
            slide_text = ' '.join(slide.get('text_blocks', []))
            for pattern, metric_type in metric_patterns:
                matches = re.finditer(pattern, slide_text)
                for match in matches:
                    context_start = max(0, match.start() - 50)
                    context_end = min(len(slide_text), match.end() + 50)
                    context = slide_text[context_start:context_end].strip()
                    if context_start > 0:
                        context = context[context.find(' ') + 1:]
                    if context_end < len(slide_text):
                        last_space = context.rfind(' ')
                        if last_space > 0:
                            context = context[:last_space]
                    metrics.append({
                        'value': match.group(0),
                        'type': metric_type,
                        'context': context,
                        'slide': slide['number'],
                    })

        seen = set()
        unique_metrics = []
        for m in metrics:
            key = (m['value'], m['slide'])
            if key not in seen:
                seen.add(key)
                unique_metrics.append(m)

        if len(unique_metrics) > 20:
            logger.info(f"  Metrics truncated: showing 20 of {len(unique_metrics)}")
        return unique_metrics[:20]

    def extract_timelines(self, slides_data: List[Dict]) -> List[Dict]:
        timelines = []
        timeline_patterns = [
            r'(Q[1-4]\s*(?:19[7-9]\d|20\d{2}))\s*[:\-–]\s*(.+?)(?:\n|$)',
            r'((?:19[7-9]\d|20\d{2}))\s*[:\-–]\s*(.+?)(?:\n|$)',
            r'((?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\w*\s+(?:19[7-9]\d|20\d{2}))\s*[:\-–]\s*(.+?)(?:\n|$)',
        ]

        for slide in slides_data:
            if slide.get('is_boilerplate'):
                continue
            slide_text = '\n'.join(slide.get('text_blocks', []))
            for pattern in timeline_patterns:
                matches = re.finditer(pattern, slide_text, re.IGNORECASE)
                for match in matches:
                    timelines.append({
                        'date': match.group(1).strip(),
                        'milestone': match.group(2).strip()[:100],
                        'slide': slide['number'],
                    })

        if len(timelines) > 15:
            logger.info(f"  Timelines truncated: showing 15 of {len(timelines)}")
        return timelines[:15]

    def calculate_quality_score(self, slides_data: List[Dict], num_slides: int,
                                figures: List[Dict], executive_summary: str,
                                key_takeaways: str, metadata: Dict,
                                has_summary: bool) -> Dict:
        scores = {}

        total_chars = sum(len(' '.join(s.get('text_blocks', []))) for s in slides_data)
        chars_per_slide = total_chars / max(num_slides, 1)
        scores['text_extraction'] = round(min(10, chars_per_slide / 60), 1)

        titled_slides = sum(1 for s in slides_data if s.get('title'))
        scores['structure_recognition'] = round((titled_slides / max(num_slides, 1)) * 10, 1)

        if figures:
            described = sum(1 for f in figures if f.get('description'))
            scores['visual_analysis'] = round(min(10, (described / max(len(figures), 1)) * 10), 1)
        else:
            scores['visual_analysis'] = 7.0

        summary_score = 0.0
        if has_summary:
            if executive_summary:
                summary_score += min(6, len(executive_summary) / 400)
            if key_takeaways:
                takeaway_count = key_takeaways.count('- ')
                summary_score += min(4, takeaway_count * 0.8)
        else:
            summary_score = 7.0
        scores['summary_quality'] = round(min(10, summary_score), 1)

        meta_fields = ['title', 'date', 'author', 'classification']
        meta_present = sum(1 for f in meta_fields if metadata.get(f))
        scores['metadata_completeness'] = round((meta_present / len(meta_fields)) * 10, 1)

        weights = {
            'text_extraction': 0.30,
            'structure_recognition': 0.20,
            'visual_analysis': 0.15,
            'summary_quality': 0.25,
            'metadata_completeness': 0.10,
        }
        overall = sum(scores[k] * weights[k] for k in weights)
        scores['overall'] = round(max(0, min(10, overall)), 1)

        if scores['overall'] >= 8.0:
            scores['assessment'] = 'Excellent'
        elif scores['overall'] >= 5.5:
            scores['assessment'] = 'Good'
        elif scores['overall'] >= 4.0:
            scores['assessment'] = 'Fair'
        else:
            scores['assessment'] = 'Poor - Consider manual review'

        return scores

    @staticmethod
    def _yaml_escape(value: str) -> str:
        return value.replace('\\', '\\\\').replace('"', '\\"')

    def generate_markdown(self, file_info: Dict, content: Dict, classification: str,
                          content_type: str, audience: Optional[str], topics: List[str],
                          figures: List[Dict], chemical_structures: List[Dict],
                          executive_summary: str, key_takeaways: str,
                          action_items: List[Dict], key_metrics: List[Dict],
                          timelines: List[Dict], quality_scores: Dict,
                          language: str = 'en') -> str:
        md = []
        slides_data = content['slides']
        has_summary = bool(executive_summary)

        # YAML frontmatter
        md.append("---")
        md.append(f"title: \"{self._yaml_escape(file_info.get('title', 'Untitled'))}\"")
        md.append(f"filename: \"{self._yaml_escape(file_info['filename'])}\"")
        if file_info.get('date'):
            md.append(f"date: \"{file_info['date']}\"")
        if content['metadata'].get('author'):
            md.append(f"author: \"{self._yaml_escape(content['metadata']['author'])}\"")
        if content['metadata'].get('last_modified_by'):
            md.append(f"last_modified_by: \"{self._yaml_escape(content['metadata']['last_modified_by'])}\"")
        if content['metadata'].get('revision'):
            md.append(f"revision: {content['metadata']['revision']}")
        md.append(f"classification: {classification}")
        if audience:
            md.append(f"audience: {audience}")
        if topics:
            md.append("topics:")
            for t in topics:
                md.append(f"  - {t}")
        md.append(f"content_type: {content_type}")
        md.append(f"language: {language}")
        md.append(f"source_format: {file_info['source_format']}")
        md.append(f"num_slides: {content['num_slides']}")
        md.append(f"has_speaker_notes: {str(content.get('has_speaker_notes', False)).lower()}")
        md.append(f"has_chemical_structures: {str(bool(chemical_structures)).lower()}")
        if chemical_structures:
            md.append(f"chemical_structures_count: {len(chemical_structures)}")
        md.append(f"has_action_items: {str(bool(action_items)).lower()}")
        md.append(f"has_summary: {str(has_summary).lower()}")
        md.append(f"extraction_method: {content['extraction_method']}")
        if not self.no_vision:
            md.append(f"vision_model: {self.VISION_MODEL}")
        md.append(f"processing_date: {datetime.now().strftime('%Y-%m-%d')}")
        md.append(f"quality_overall: {quality_scores['overall']}/10")
        md.append(f"quality_assessment: {quality_scores['assessment']}")
        md.append("---\n")

        # Title and header
        md.append(f"# {file_info.get('title', 'Untitled')}\n")
        header_parts = []
        if file_info.get('date'):
            header_parts.append(f"**Date**: {file_info['date']}")
        if content['metadata'].get('author'):
            header_parts.append(f"**Author**: {content['metadata']['author']}")
        header_parts.append(f"**Classification**: {classification.upper()}")
        md.append("  |  ".join(header_parts))

        secondary_parts = []
        if audience:
            audience_display = audience.replace('_', ' ').title()
            secondary_parts.append(f"**Audience**: {audience_display}")
        if topics:
            secondary_parts.append(f"**Topics**: {', '.join(topics)}")
        if secondary_parts:
            md.append("  |  ".join(secondary_parts))
        md.append("\n---\n")

        # Executive Summary
        if has_summary:
            md.append("## Executive Summary\n")
            md.append(f"{executive_summary}\n")
            md.append("---\n")

        # Slide Content
        md.append("## Slide Content\n")
        recurring_footers = self._detect_recurring_footer(slides_data) or []
        for slide in slides_data:
            if slide.get('is_boilerplate'):
                continue
            title = slide.get('title') or f"Slide {slide['number']}"
            md.append(f"### Slide {slide['number']}: {title}\n")

            slide_num_str = str(slide['number'])
            blocks = [b for b in slide.get('text_blocks', [])
                      if b != title
                      and b.strip() != slide_num_str
                      and b.strip() not in recurring_footers]
            blocks = self._merge_short_fragments(blocks)
            for block in blocks:
                md.append(f"{block}\n")

            for table in slide.get('tables', []):
                if table:
                    md.append("")
                    header = table[0]
                    md.append("| " + " | ".join(str(c) for c in header) + " |")
                    md.append("|" + "|".join("---" for _ in header) + "|")
                    for row in table[1:]:
                        md.append("| " + " | ".join(str(c) for c in row) + " |")
                    md.append("")

            if slide.get('notes'):
                md.append(f"**Notes**: {slide['notes']}\n")

            if slide.get('visual_elements'):
                md.append(f"**Visual elements**: {slide['visual_elements']}\n")

            md.append("")

        md.append("---\n")

        # Figures & Charts
        if figures:
            md.append("## Figures & Charts\n")
            for i, fig in enumerate(figures, 1):
                title = fig.get('title', fig.get('label', f'Figure {i}'))
                md.append(f"### Figure {i}: {title} (Slide {fig.get('slide', '?')})\n")
                md.append(f"**Type**: {fig.get('type', 'unknown')}")
                if fig.get('description'):
                    md.append(f"**Description**: {fig['description']}")
                if fig.get('x_axis'):
                    md.append(f"**X-axis**: {fig['x_axis']}  |  **Y-axis**: {fig.get('y_axis', 'N/A')}")
                if fig.get('key_findings'):
                    md.append(f"**Key Findings**: {fig['key_findings']}")
                md.append("")
            md.append("---\n")

        # Chemical Structures
        if chemical_structures:
            md.append("## Chemical Structures\n")
            for i, struct in enumerate(chemical_structures, 1):
                label = struct.get('label', f'Structure {i}')
                md.append(f"### Structure {i}: {label} (Slide {struct.get('slide', '?')})\n")
                if struct.get('smiles'):
                    conf = struct.get('smiles_confidence', 'unknown')
                    md.append(f"**SMILES** (confidence: {conf}): `{struct['smiles']}`")
                if struct.get('molecular_formula'):
                    mw = struct.get('molecular_weight', 'N/A')
                    md.append(f"**Molecular Formula**: {struct['molecular_formula']}  |  **MW**: {mw}")
                if struct.get('description'):
                    md.append(f"**Description**: {struct['description']}")
                md.append("")
            md.append("---\n")

        # Key Metrics & Data Points
        if key_metrics:
            md.append("## Key Metrics & Data Points\n")
            for m in key_metrics:
                md.append(f"- {m['context']} (Slide {m['slide']})")
            md.append("\n---\n")

        # Action Items
        if action_items:
            md.append("## Action Items\n")
            md.append("| Action | Owner | Deadline | Slide |")
            md.append("|--------|-------|----------|-------|")
            for item in action_items:
                owner = item.get('owner') or 'TBD'
                deadline = item.get('deadline') or 'TBD'
                md.append(f"| {item['action'][:80]} | {owner} | {deadline} | {item['slide']} |")
            md.append("\n---\n")

        # Timeline & Milestones
        if timelines:
            md.append("## Timeline & Milestones\n")
            for t in timelines:
                md.append(f"- **{t['date']}**: {t['milestone']}")
            md.append("\n---\n")

        # References & Links
        hyperlinks = content.get('hyperlinks', [])
        media_flags = content.get('media_flags', [])
        if hyperlinks or media_flags:
            md.append("## References & Links\n")
            for link in hyperlinks:
                md.append(f"- [{link['text'] or link['url']}]({link['url']}) (Slide {link['slide']})")
            for media in media_flags:
                md.append(f"- Slide {media['slide']}: embedded {media['type']} (not extractable)")
            md.append("\n---\n")

        # Key Takeaways
        if key_takeaways:
            md.append("## Key Takeaways\n")
            md.append(f"{key_takeaways}\n")
            md.append("---\n")

        # Quality Assessment
        md.append("## Quality Assessment\n")
        md.append(f"**Overall**: {quality_scores['overall']}/10 — {quality_scores['assessment']}\n")
        md.append("**Component Scores**:")
        md.append(f"- Text extraction: {quality_scores['text_extraction']}/10 (weight 0.30)")
        md.append(f"- Structure recognition: {quality_scores['structure_recognition']}/10 (weight 0.20)")
        md.append(f"- Visual analysis: {quality_scores['visual_analysis']}/10 (weight 0.15)")
        md.append(f"- Summary quality: {quality_scores['summary_quality']}/10 (weight 0.25)")
        md.append(f"- Metadata completeness: {quality_scores['metadata_completeness']}/10 (weight 0.10)")

        return "\n".join(md)

    def generate_output_filename(self, file_info: Dict, classification: str = 'confidential') -> str:
        title_clean = re.sub(r'[^a-zA-Z0-9\s]', '', file_info.get('title', 'untitled'))
        title_words = title_clean.split()[:5]
        title_slug = '_'.join(w.lower() for w in title_words) or 'untitled'
        date = file_info.get('date') or 'undated'

        template = self.NAMING_SCHEMES.get(self.naming, self.NAMING_SCHEMES['default'])
        filename = template.format(
            title_slug=title_slug,
            date=date,
            classification=classification,
        )
        return sanitize_filename(f"{filename}.md")

    def check_if_processed(self, file_info: Dict, classification: str = 'confidential') -> bool:
        expected = self.generate_output_filename(file_info, classification)
        return expected.replace('.md', '') in self._existing_files

    def parse_json_response(self, text: str) -> Any:
        text = text.strip()
        code_block = re.search(r'```(?:json)?\s*\n?(.*?)\n?```', text, re.DOTALL)
        if code_block:
            text = code_block.group(1).strip()

        try:
            return json.loads(text)
        except json.JSONDecodeError:
            pass

        array_match = re.search(r'(\[.*\])', text, re.DOTALL)
        if array_match:
            try:
                return json.loads(array_match.group(1))
            except json.JSONDecodeError:
                pass

        obj_match = re.search(r'(\{.*\})', text, re.DOTALL)
        if obj_match:
            try:
                return json.loads(obj_match.group(1))
            except json.JSONDecodeError:
                pass

        return None

    def process_single_presentation(self, filepath: Path, skip_existing: bool = True) -> bool:
        logger.info(f"\n{'='*60}")
        logger.info(f"Processing: {filepath.name}")
        logger.info(f"{'='*60}")

        # File size check
        file_size_mb = filepath.stat().st_size / (1024 * 1024)
        if file_size_mb > self.MAX_FILE_SIZE_MB:
            logger.error(f"  File too large: {file_size_mb:.1f}MB (max: {self.MAX_FILE_SIZE_MB}MB) — skipping")
            return False

        # Step 1: Parse filename
        file_info = self.parse_filename(filepath)
        logger.info(f"  Title: {file_info['title']}")
        logger.info(f"  Format: {file_info['source_format']}")
        logger.info(f"  Date: {file_info.get('date', 'unknown')}")

        # Early classification from filename only (for skip-check and naming)
        filename_classification, _ = self.detect_classification(file_info['filename'], '')

        if skip_existing and self.check_if_processed(file_info, filename_classification):
            logger.info(f"  Already processed — skipping")
            return True

        # Step 2: Extract content
        logger.info("\n--- Extracting content ---")
        if file_info['source_format'] == 'pptx':
            content = self.extract_pptx_content(filepath)
        else:
            content = self.extract_pdf_content(filepath)

        if not content['slides']:
            logger.error("Failed to extract content")
            self.log_processing(file_info['filename'], file_info['source_format'], 0,
                              'unknown', 'unknown', 0, None, "FAILED_EXTRACTION")
            return False

        num_slides = content['num_slides']
        total_chars = sum(len(' '.join(s.get('text_blocks', []))) for s in content['slides'])
        logger.info(f"  Extracted {num_slides} slides, {total_chars} chars")

        # Detect language
        language = self._detect_language(content['slides'])
        if language != 'en':
            logger.info(f"  Language detected: {language}")

        # Step 3: Full classification (filename + content)
        logger.info("\n--- Detecting classification ---")
        all_text = ' '.join(' '.join(s.get('text_blocks', [])) for s in content['slides'])
        classification, classification_signals = self.detect_classification(file_info['filename'], all_text)
        logger.info(f"  Classification: {classification}")
        for signal in classification_signals:
            logger.info(f"    Signal: {signal}")

        # Step 4: Content type classification
        content_type = self.classify_content_type(content['slides'], file_info['filename'])
        logger.info(f"  Content type: {content_type}")

        # Step 5: Vision AI analysis (figures, audience, topics, structures)
        # Gate: only run if there's visual content that text extraction can't capture
        run_vision_analysis = not self.no_vision
        if run_vision_analysis and file_info['source_format'] == 'pptx':
            if not content.get('has_visual_content'):
                run_vision_analysis = False
                logger.info("\n--- Vision AI analysis: skipped (no images/charts detected in PPTX) ---")
            else:
                logger.info(f"\n--- Vision AI analysis ({content.get('image_count', 0)} images, "
                           f"{content.get('chart_count', 0)} charts) ---")
        else:
            logger.info("\n--- Vision AI analysis ---")

        # Convert images once and reuse for vision analysis + structure refinement
        slide_images = None
        if run_vision_analysis:
            if file_info['source_format'] == 'pptx':
                slide_images = self._convert_pptx_to_images(filepath)
            else:
                slide_images = self._convert_pdf_to_images(filepath)

        vision_result = self.analyze_with_vision(filepath, content['slides'], file_info['source_format'],
                                                  skip=not run_vision_analysis,
                                                  preloaded_images=slide_images)

        if vision_result.get('content_type_hint'):
            content_type = vision_result['content_type_hint']
            logger.info(f"  Content type (Vision AI): {content_type}")

        audience = vision_result.get('audience')
        topics = vision_result.get('topics', [])
        figures = vision_result.get('figures', [])
        chemical_structures = vision_result.get('chemical_structures', [])

        if vision_result.get('classification_visual_cues'):
            for cue in vision_result['classification_visual_cues']:
                classification_signals.append(f"visual: {cue}")

        # Step 6: Chemical structure refinement
        if chemical_structures and not file_info.get('no_structures'):
            logger.info(f"\n--- Chemical structure refinement ({len(chemical_structures)} structures) ---")
            chemical_structures = self.refine_chemical_structures(
                filepath, chemical_structures, content['slides'], file_info['source_format'],
                preloaded_images=slide_images)
        elif file_info.get('no_structures'):
            logger.info("  Skipping SMILES refinement ([No Structures] in filename)")
            chemical_structures = []

        # Step 7: Summary generation (conditional)
        executive_summary = ''
        key_takeaways = ''
        if content_type not in ('agenda', 'template'):
            logger.info("\n--- Generating summary ---")
            executive_summary, key_takeaways = self.generate_summary_and_takeaways(
                content['slides'], file_info, language=language)
            if executive_summary:
                logger.info(f"  Summary: {len(executive_summary)} chars")
        else:
            logger.info(f"  Skipping summary (content_type={content_type})")

        # Step 8: Extract action items, metrics, timelines
        action_items = self.extract_action_items(content['slides'])
        key_metrics = self.extract_key_metrics(content['slides'])
        timelines = self.extract_timelines(content['slides'])
        logger.info(f"  Action items: {len(action_items)}, Metrics: {len(key_metrics)}, Timelines: {len(timelines)}")

        # Step 9: Quality scoring
        logger.info("\n--- Quality Assessment ---")
        metadata_for_quality = {
            'title': file_info.get('title'),
            'date': file_info.get('date'),
            'author': content['metadata'].get('author'),
            'classification': classification,
        }
        quality_scores = self.calculate_quality_score(
            slides_data=content['slides'],
            num_slides=num_slides,
            figures=figures + chemical_structures,
            executive_summary=executive_summary,
            key_takeaways=key_takeaways,
            metadata=metadata_for_quality,
            has_summary=bool(executive_summary),
        )
        logger.info(f"  Overall: {quality_scores['overall']}/10 — {quality_scores['assessment']}")

        # Step 10: Generate markdown
        logger.info("\n--- Generating markdown ---")
        markdown = self.generate_markdown(
            file_info=file_info,
            content=content,
            classification=classification,
            content_type=content_type,
            audience=audience,
            topics=topics,
            figures=figures,
            chemical_structures=chemical_structures,
            executive_summary=executive_summary,
            key_takeaways=key_takeaways,
            action_items=action_items,
            key_metrics=key_metrics,
            timelines=timelines,
            quality_scores=quality_scores,
            language=language,
        )

        # Save output
        output_filename = self.generate_output_filename(file_info, classification)
        output_path = self.output_dir / output_filename
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown)
        self._existing_files.add(output_path.stem)

        self.log_processing(file_info['filename'], file_info['source_format'], num_slides,
                           classification, content_type, total_chars, quality_scores, "SAVED")

        logger.info(f"\n{'='*60}")
        logger.info(f"SAVED: {output_path}")
        logger.info(f"  Slides: {num_slides} | Chars: {total_chars}")
        logger.info(f"  Classification: {classification} | Type: {content_type}")
        logger.info(f"  Quality: {quality_scores['overall']}/10 — {quality_scores['assessment']}")
        logger.info(f"{'='*60}")

        return True

    def process_all_presentations(self, skip_existing: bool = True):
        logger.info("=" * 70)
        logger.info("PRESENTATION PROCESSING PIPELINE")
        logger.info("=" * 70)
        logger.info(f"Input folder: {self.input_folder}")
        logger.info(f"Output directory: {self.output_dir}")
        logger.info(f"Recursive: {self.recursive}")
        logger.info(f"Vision AI: {'disabled' if self.no_vision else 'enabled'}")
        logger.info(f"Vision model: {self.VISION_MODEL}")
        logger.info("=" * 70)

        files = self.list_presentation_files()
        if not files:
            logger.error("No presentation files found")
            return

        success_count = 0
        fail_count = 0

        for i, filepath in enumerate(files):
            logger.info(f"\n[{i+1}/{len(files)}] {filepath.name}")
            try:
                if self.process_single_presentation(filepath, skip_existing=skip_existing):
                    success_count += 1
                else:
                    fail_count += 1
            except Exception as e:
                logger.error(f"Unexpected error processing {filepath.name}: {e}")
                fail_count += 1

        logger.info(f"\n{'='*70}")
        logger.info(f"PIPELINE COMPLETE: {success_count} succeeded, {fail_count} failed "
                   f"out of {len(files)}")
        logger.info(f"Processing log: {self.processing_log_path}")
        logger.info(f"{'='*70}")


def main():
    parser = argparse.ArgumentParser(
        description="Process presentations (PPTX/PDF) into structured markdown"
    )
    parser.add_argument("--input", default="test_presentation",
                       help="Folder containing presentations (PPTX/PDF)")
    parser.add_argument("--output", default="output_presentations",
                       help="Output directory for markdown files")
    parser.add_argument("--single", default=None,
                       help="Process a single file only")
    parser.add_argument("--recursive", action="store_true",
                       help="Also scan subfolders for presentations")
    parser.add_argument("--no-skip", action="store_true",
                       help="Reprocess files that already exist")
    parser.add_argument("--no-vision", action="store_true",
                       help="Skip Vision AI (text-only extraction, no API calls)")
    parser.add_argument("--naming", default="default",
                       choices=["default", "dated", "classified"],
                       help="Output filename scheme: "
                            "'default' = presentation_{title}; "
                            "'dated' = {date}_{classification}_{title}; "
                            "'classified' = {classification}_{date}_{title}")
    parser.add_argument("--verbose", action="store_true",
                       help="Enable debug logging")

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    # Validate paths
    try:
        validate_path(args.input, must_exist=True, allow_dir=True, allow_file=False)
        validate_output_path(args.output)
        if args.single:
            validate_path(args.single, must_exist=True, allow_file=True, allow_dir=False)
    except (ValueError, FileNotFoundError) as e:
        logger.error(f"Path validation failed: {e}")
        return

    pipeline = PresentationPipeline(
        input_folder=args.input,
        output_dir=args.output,
        no_vision=args.no_vision,
        recursive=args.recursive,
        naming=args.naming,
    )

    if args.single:
        single_path = Path(args.single)
        pipeline.process_single_presentation(single_path, skip_existing=not args.no_skip)
    else:
        pipeline.process_all_presentations(skip_existing=not args.no_skip)


if __name__ == "__main__":
    main()
