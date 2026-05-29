"""
Competitive Intelligence Processing Pipeline
Converts competitive intelligence documents (PDFs, PPTXs) to structured markdown
organized by company and drug asset.

Supports ONCrg reports (Sentinel, Conference, Pipeline Strategy), Beacon reports,
industry analyst reports (Frost & Sullivan, Catenion, Clarivate, etc.), and
scientific papers. Uses Claude Sonnet for entity extraction (companies, assets,
clinical events, deals).
"""

import os
import re
import json
import argparse
import time
import threading
from pathlib import Path
from typing import List, Dict, Tuple, Optional, Any
from concurrent.futures import ThreadPoolExecutor, as_completed
import logging
from datetime import datetime

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

from auth import get_anthropic_client
from pipeline_security import validate_path, validate_output_path, sanitize_filename, check_pptx_safe


class CIPipeline:
    """Pipeline for processing competitive intelligence documents into structured markdown."""

    VISION_MODEL = "claude-sonnet-4-6"
    MAX_FILE_SIZE_MB = 500

    DEFAULT_BUDGET_CAP = 30
    API_MAX_RETRIES = 3
    API_RETRY_DELAY = 2.0
    MAX_CONCURRENT_CHUNKS = 4

    CHUNK_SIZE_CHARS = 12000
    CHUNK_OVERLAP_CHARS = 500

    SUPPORTED_EXTENSIONS = {'.pdf', '.pptx'}

    EXTRACTION_SYSTEM_PROMPT = """You extract pharmaceutical competitive intelligence data from ADC (Antibody-Drug Conjugate) document sections.
Return ONLY valid JSON with this structure:
{"companies":{"Company Name":{"assets":[{"name":"drug name or code","target":"molecular target or null","payload":"payload type or null","modality":"ADC type or null","indications":["cancer types"],"stage":"Phase I/II/III/Approved/Preclinical or null","key_updates":["concise updates"],"clinical_data":["efficacy/safety data: ORR, PFS, OS, AEs"]}]}},"deals":[{"parties":["Company A","Company B"],"type":"acquisition/licensing/collaboration","details":"brief description"}],"regulatory_events":[{"company":"Name","asset":"drug","event":"FDA approval/Priority Review/IND/filing/etc."}]}
Rules: Use canonical company names. Include ALL companies mentioned. Capture drug names exactly as written. If clinical data is mentioned, capture it in clinical_data. If nothing found, return {"companies":{},"deals":[],"regulatory_events":[]}."""

    ASSET_CODE_RE = re.compile(r'^[A-Z]{1,5}[-]?\d{2,}', re.IGNORECASE)

    DOC_TYPE_PATTERNS = {
        'sentinel': re.compile(r'ONCrgSentinel', re.IGNORECASE),
        'conference': re.compile(r'ONCrgConference', re.IGNORECASE),
        'strategy': re.compile(r'ONCrgPipeStrat', re.IGNORECASE),
        'beacon': re.compile(r'^Beacon', re.IGNORECASE),
        'industry_report': re.compile(
            r'^(Frost|Catenion|Clarivate|GlobalData|Morgan|Guggenheim|Scitaris|citeline|jpm|IDEA)',
            re.IGNORECASE
        ),
        'paper': re.compile(
            r'(\bet\s+al\b|20\d{6}\b)',
            re.IGNORECASE
        ),
    }

    SOURCE_MAP = {
        'sentinel': 'ONCrg',
        'conference': 'ONCrg',
        'strategy': 'ONCrg',
        'beacon': 'Beacon',
    }

    INDUSTRY_SOURCE_PATTERNS = [
        (re.compile(r'^Frost', re.IGNORECASE), 'Frost & Sullivan'),
        (re.compile(r'^Catenion', re.IGNORECASE), 'Catenion'),
        (re.compile(r'^Clarivate', re.IGNORECASE), 'Clarivate'),
        (re.compile(r'^GlobalData', re.IGNORECASE), 'GlobalData'),
        (re.compile(r'^Morgan', re.IGNORECASE), 'Morgan Stanley'),
        (re.compile(r'^Guggenheim', re.IGNORECASE), 'Guggenheim'),
        (re.compile(r'^Scitaris', re.IGNORECASE), 'Scitaris'),
        (re.compile(r'^citeline', re.IGNORECASE), 'Citeline'),
        (re.compile(r'^jpm', re.IGNORECASE), 'JP Morgan'),
        (re.compile(r'^IDEA', re.IGNORECASE), 'IDEA Pharma'),
        (re.compile(r'^AdisInsight', re.IGNORECASE), 'AdisInsight'),
        (re.compile(r'^BioCentury', re.IGNORECASE), 'BioCentury'),
        (re.compile(r'^QPS', re.IGNORECASE), 'QPS'),
    ]

    MONTH_MAP = {
        'january': 1, 'february': 2, 'march': 3, 'april': 4, 'may': 5, 'june': 6,
        'july': 7, 'august': 8, 'september': 9, 'october': 10, 'november': 11, 'december': 12,
        'jan': 1, 'feb': 2, 'mar': 3, 'apr': 4, 'may': 5, 'jun': 6,
        'jul': 7, 'aug': 8, 'sep': 9, 'oct': 10, 'nov': 11, 'dec': 12,
    }

    NAMING_SCHEMES = {
        'default': 'ci_{source_type}_{title_slug}',
        'dated': 'ci_{date}_{source_type}_{title_slug}',
        'source': 'ci_{source}_{date}_{title_slug}',
    }

    def __init__(self, input_folder: str, output_dir: str = "output_ci",
                 recursive: bool = False, naming: str = 'default',
                 budget: int = 30):
        self.input_folder = Path(input_folder)
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.recursive = recursive
        self.naming = naming

        self._client = None
        self._budget_cap = budget if budget > 0 else None
        self._api_call_count = 0
        self._api_lock = threading.Lock()

        self._existing_files = set(f.stem for f in self.output_dir.glob('*.md'))

        self.processing_log_path = self.output_dir / "processing_log.tsv"
        self._init_processing_log()

    # ─── Infrastructure ──────────────────────────────────────────────────────

    def _init_processing_log(self):
        if not self.processing_log_path.exists() or self.processing_log_path.stat().st_size == 0:
            with open(self.processing_log_path, 'w', encoding='utf-8') as f:
                f.write("TIMESTAMP\tFILENAME\tDOC_TYPE\tSOURCE\tPAGES\tCHARS\t"
                        "COMPANIES\tASSETS\tQUALITY\tSTATUS\n")

    def _log_processing(self, filename: str, doc_type: str, source: str,
                        pages: int, chars: int, companies: int, assets: int,
                        quality: float, status: str):
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        with open(self.processing_log_path, 'a', encoding='utf-8') as f:
            f.write(f"{timestamp}\t{filename}\t{doc_type}\t{source}\t{pages}\t"
                    f"{chars}\t{companies}\t{assets}\t{quality}\t{status}\n")

    def _get_client(self):
        if self._client is None:
            self._client = get_anthropic_client()
        return self._client

    def _track_api_call(self, count: int = 1) -> bool:
        with self._api_lock:
            if self._budget_cap and (self._api_call_count + count) > self._budget_cap:
                logger.warning(f"  Budget limit reached: {self._api_call_count}/{self._budget_cap}")
                return False
            self._api_call_count += count
            return True

    def _api_call_with_retry(self, call_fn, description: str = "API call"):
        delay = self.API_RETRY_DELAY
        for attempt in range(self.API_MAX_RETRIES):
            try:
                return call_fn()
            except Exception as e:
                error_str = str(e)
                is_transient = any(s in error_str for s in
                                   ('429', '503', '529', 'timeout', 'Timeout', 'overloaded'))
                if is_transient and attempt < self.API_MAX_RETRIES - 1:
                    logger.warning(f"  {description} attempt {attempt + 1} failed, retrying in {delay:.0f}s...")
                    time.sleep(delay)
                    delay *= 2
                else:
                    raise
        return None

    def parse_json_response(self, text: str) -> Optional[Any]:
        text = text.strip()
        code_block = re.search(r'```(?:json)?\s*\n?(.*?)\n?```', text, re.DOTALL)
        if code_block:
            text = code_block.group(1).strip()
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            pass
        for start_char, end_char in [('[', ']'), ('{', '}')]:
            start = text.find(start_char)
            if start >= 0:
                depth = 0
                for i, c in enumerate(text[start:], start):
                    if c == start_char:
                        depth += 1
                    elif c == end_char:
                        depth -= 1
                        if depth == 0:
                            try:
                                return json.loads(text[start:i + 1])
                            except json.JSONDecodeError:
                                break
        return None

    @staticmethod
    def _escape_yaml(text: str) -> str:
        if not text:
            return ''
        return text.replace('\\', '\\\\').replace('"', '\\"').replace('\n', ' ')

    # ─── File Discovery ──────────────────────────────────────────────────────

    def list_ci_files(self) -> List[Path]:
        files = []
        for ext in ('*.pdf', '*.pptx'):
            if self.recursive:
                files.extend(self.input_folder.rglob(ext))
            else:
                files.extend(self.input_folder.glob(ext))
        files = sorted(files, key=lambda p: p.name.lower())
        pdf_count = sum(1 for f in files if f.suffix.lower() == '.pdf')
        pptx_count = sum(1 for f in files if f.suffix.lower() == '.pptx')
        logger.info(f"Found {len(files)} CI documents ({pdf_count} PDF, {pptx_count} PPTX)")
        return files

    # ─── Document Type Detection ─────────────────────────────────────────────

    def detect_document_type(self, filepath: Path) -> Tuple[str, str]:
        """Classify document type and source from filename."""
        stem = filepath.stem

        for doc_type, pattern in self.DOC_TYPE_PATTERNS.items():
            if pattern.search(stem):
                source = self._resolve_source(stem, doc_type)
                return doc_type, source

        return 'other', 'Unknown'

    def _resolve_source(self, stem: str, doc_type: str) -> str:
        if doc_type in self.SOURCE_MAP:
            return self.SOURCE_MAP[doc_type]
        if doc_type == 'industry_report':
            for pattern, name in self.INDUSTRY_SOURCE_PATTERNS:
                if pattern.search(stem):
                    return name
            return 'Industry'
        if doc_type == 'paper':
            return 'Academic'
        return 'Unknown'

    def _extract_date_from_filename(self, filepath: Path) -> Optional[str]:
        """Extract date from CI document filename."""
        stem = filepath.stem

        # ONCrgSentinel-ADC-Month-Year pattern
        month_year = re.search(
            r'(January|February|March|April|May|June|July|August|September|'
            r'October|November|December|Jan|Feb|Mar|Apr|Jun|Jul|Aug|Sep|Oct|Nov|Dec)'
            r'[-_\s]*(\d{4})',
            stem, re.IGNORECASE
        )
        if month_year:
            month_str = month_year.group(1).lower()
            year = month_year.group(2)
            month = self.MONTH_MAP.get(month_str, 1)
            return f"{year}-{month:02d}-01"

        # Quarter-Year pattern (Q1-2025)
        quarter_match = re.search(r'Q([1-4])[-_\s]*(\d{4})', stem)
        if quarter_match:
            q = int(quarter_match.group(1))
            year = quarter_match.group(2)
            month = (q - 1) * 3 + 1
            return f"{year}-{month:02d}-01"

        # Trailing year (2024, 2025, 2026)
        year_match = re.search(r'[-_\s](\d{4})(?:[-_.\s]|$)', stem)
        if year_match:
            return f"{year_match.group(1)}-01-01"

        # YYYYMMDD prefix in filename
        prefix_match = re.search(r'(\d{4})(\d{2})(\d{2})', stem)
        if prefix_match:
            return f"{prefix_match.group(1)}-{prefix_match.group(2)}-{prefix_match.group(3)}"

        # Two-digit year-month codes like 2303, 2307, 2404
        code_match = re.search(r'[-_\s](\d{2})(\d{2})(?:[-_.\s]|$)', stem)
        if code_match:
            yy = int(code_match.group(1))
            mm = int(code_match.group(2))
            if 20 <= yy <= 30 and 1 <= mm <= 12:
                return f"20{yy}-{mm:02d}-01"

        return None

    def _extract_title_from_filename(self, filepath: Path, doc_type: str) -> str:
        """Extract a human-readable title from the filename."""
        stem = filepath.stem

        if doc_type == 'sentinel':
            match = re.search(r'ONCrgSentinel-ADC-(.+)', stem)
            if match:
                return f"ONCrg Sentinel ADC — {match.group(1).replace('-', ' ')}"
            return f"ONCrg Sentinel — {stem}"

        if doc_type == 'conference':
            match = re.search(r'ONCrgConference-ADC-(.+)', stem)
            if match:
                parts = match.group(1).replace('-', ' ')
                return f"ONCrg Conference ADC — {parts}"
            return f"ONCrg Conference — {stem}"

        if doc_type == 'strategy':
            match = re.search(r'ONCrgPipeStrategies-ADC-Presentation-(.+)', stem)
            if match:
                parts = match.group(1).replace('-', ' ')
                return f"ONCrg Pipeline Strategies ADC — {parts}"
            return f"ONCrg Pipeline Strategies — {stem}"

        # General: clean up the stem
        title = re.sub(r'[-_]+', ' ', stem)
        title = re.sub(r'\s+', ' ', title).strip()
        return title

    # ─── Content Extraction ──────────────────────────────────────────────────

    def extract_pdf_content(self, pdf_path: Path) -> Dict:
        """Extract text and tables from a PDF."""
        import fitz
        import pdfplumber

        pages = []
        total_chars = 0

        try:
            doc = fitz.open(str(pdf_path))
            num_pages = len(doc)

            for page_idx in range(num_pages):
                page = doc[page_idx]
                text = page.get_text()
                char_count = len(text.strip())
                total_chars += char_count
                pages.append({
                    'page_num': page_idx + 1,
                    'text': text,
                    'char_count': char_count,
                    'tables': [],
                })

            doc.close()

            # Table extraction with pdfplumber
            with pdfplumber.open(str(pdf_path)) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    extracted_tables = page.extract_tables()
                    if extracted_tables and page_idx < len(pages):
                        for table in extracted_tables:
                            cleaned = [[cell or '' for cell in row] for row in table if row]
                            if cleaned and len(cleaned) >= 2:
                                pages[page_idx]['tables'].append(cleaned)

        except Exception as e:
            logger.error(f"  PDF extraction failed: {e}")
            return {'pages': [], 'full_text': '', 'total_pages': 0,
                    'total_chars': 0, 'extraction_method': 'failed'}

        full_text = '\n\n'.join(p['text'].strip() for p in pages if p['text'].strip())

        return {
            'pages': pages,
            'full_text': full_text,
            'total_pages': num_pages,
            'total_chars': total_chars,
            'extraction_method': 'native_text',
        }

    def extract_pptx_content(self, pptx_path: Path) -> Dict:
        """Extract text and tables from a PPTX presentation."""
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        if not check_pptx_safe(pptx_path):
            logger.error(f"  PPTX safety check failed: {pptx_path.name}")
            return {'pages': [], 'full_text': '', 'total_pages': 0,
                    'total_chars': 0, 'extraction_method': 'failed'}

        prs = Presentation(str(pptx_path))
        pages = []
        total_chars = 0

        for slide_idx, slide in enumerate(prs.slides, 1):
            text_blocks = []
            tables = []

            if slide.shapes.title and slide.shapes.title.text.strip():
                text_blocks.append(slide.shapes.title.text.strip())

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    if table_data and len(table_data) >= 2:
                        tables.append(table_data)
                elif shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            text_blocks.append(para_text)

            slide_text = '\n'.join(text_blocks)
            char_count = len(slide_text)
            total_chars += char_count

            pages.append({
                'page_num': slide_idx,
                'text': slide_text,
                'char_count': char_count,
                'tables': tables,
            })

        full_text = '\n\n'.join(p['text'].strip() for p in pages if p['text'].strip())

        return {
            'pages': pages,
            'full_text': full_text,
            'total_pages': len(pages),
            'total_chars': total_chars,
            'extraction_method': 'native_text',
        }

    # ─── Text Pre-processing ────────────────────────────────────────────────

    _BOILERPLATE_PATTERNS = [
        re.compile(r'^.{5,}\.{3,}\s*\d+\s*$', re.MULTILINE),  # TOC dot leaders
        re.compile(r'\bPage\s+\d+\s+of\s+\d+\b', re.IGNORECASE),
        re.compile(r'(?i)^\s*copyright\s+\d{4}.{0,200}$', re.MULTILINE),
        re.compile(r'(?i)^\s*all\s+rights\s+reserved.{0,100}$', re.MULTILINE),
        re.compile(r'(?i)^\s*(?:confidential|proprietary)\s*[.\s]*$', re.MULTILINE),
        re.compile(r'(?i)^\s*inquiries\s+to:.{0,80}$', re.MULTILINE),
    ]

    def _preprocess_text(self, text: str) -> str:
        """Strip boilerplate content that wastes API tokens."""
        for pattern in self._BOILERPLATE_PATTERNS:
            text = pattern.sub('', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    # ─── Company & Asset Extraction (Core CI Logic) ──────────────────────────

    def _chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping chunks at paragraph boundaries."""
        paragraphs = text.split('\n\n')
        chunks = []
        current_chunk = []
        current_size = 0

        for para in paragraphs:
            para_len = len(para)
            if current_size + para_len > self.CHUNK_SIZE_CHARS and current_chunk:
                chunks.append('\n\n'.join(current_chunk))
                overlap_paras = []
                overlap_size = 0
                for p in reversed(current_chunk):
                    if overlap_size + len(p) > self.CHUNK_OVERLAP_CHARS:
                        break
                    overlap_paras.append(p)
                    overlap_size += len(p)
                overlap_paras.reverse()
                current_chunk = overlap_paras
                current_size = overlap_size

            current_chunk.append(para)
            current_size += para_len

        if current_chunk:
            chunks.append('\n\n'.join(current_chunk))

        return chunks

    CONTEXT_HINTS = {
        'sentinel': "monthly CI update (new trials, product/regulatory/company news)",
        'conference': "medical conference summary (clinical data by cancer indication)",
        'strategy': "pipeline strategy analysis (company-by-company asset tracking)",
        'beacon': "market research report on ADC landscape",
        'industry_report': "industry analysis report on ADCs",
        'paper': "scientific research paper on ADC therapeutics",
        'other': "competitive intelligence document on ADCs",
    }

    def _build_extraction_prompt(self, chunk_text: str, doc_type: str,
                                  source: str, chunk_idx: int,
                                  total_chunks: int) -> str:
        """Build a compact user-message prompt (schema is in the system prompt)."""
        hint = self.CONTEXT_HINTS.get(doc_type, self.CONTEXT_HINTS['other'])
        return f"Section {chunk_idx + 1}/{total_chunks} of a {hint} from {source}.\n\n{chunk_text}"

    def _extract_single_chunk(self, chunk: str, doc_type: str, source: str,
                               chunk_idx: int, total_chunks: int) -> Optional[Dict]:
        """Process one chunk via the API. Called from threads."""
        if not self._track_api_call():
            return None

        prompt = self._build_extraction_prompt(chunk, doc_type, source, chunk_idx, total_chunks)
        try:
            client = self._get_client()
            response = self._api_call_with_retry(
                lambda: client.messages.create(
                    model=self.VISION_MODEL,
                    system=self.EXTRACTION_SYSTEM_PROMPT,
                    max_tokens=4096,
                    messages=[{"role": "user", "content": prompt}]
                ),
                description=f"CI extraction chunk {chunk_idx + 1}/{total_chunks}"
            )
            result = self.parse_json_response(response.content[0].text)
            if not result:
                logger.warning(f"  Chunk {chunk_idx + 1}: failed to parse JSON response")
            return result
        except Exception as e:
            logger.error(f"  Chunk {chunk_idx + 1} extraction failed: {e}")
            return None

    def extract_companies_and_assets(self, full_text: str, doc_type: str,
                                      source: str) -> Dict:
        """Extract structured company/asset intelligence from document text."""
        cleaned = self._preprocess_text(full_text)
        chunks = self._chunk_text(cleaned)
        logger.info(f"  Text chunked into {len(chunks)} segments for extraction")

        workers = min(self.MAX_CONCURRENT_CHUNKS, len(chunks))
        results = [None] * len(chunks)

        with ThreadPoolExecutor(max_workers=workers) as executor:
            future_to_idx = {}
            for i, chunk in enumerate(chunks):
                future = executor.submit(
                    self._extract_single_chunk, chunk, doc_type, source, i, len(chunks))
                future_to_idx[future] = i

            for future in as_completed(future_to_idx):
                idx = future_to_idx[future]
                try:
                    results[idx] = future.result()
                except Exception as e:
                    logger.error(f"  Chunk {idx + 1} thread failed: {e}")

        chunk_results = [r for r in results if r]
        if not chunk_results:
            return {'companies': {}, 'deals': [], 'regulatory_events': [],
                    'companies_mentioned': [], 'assets_mentioned': []}

        return self._merge_ci_extractions(chunk_results)

    def _merge_ci_extractions(self, chunk_results: List[Dict]) -> Dict:
        """Merge company/asset extractions from multiple chunks."""
        merged_companies = {}
        merged_deals = []
        merged_regulatory = []

        for result in chunk_results:
            if not isinstance(result, dict):
                continue
            # Merge companies
            companies_raw = result.get('companies', {})
            if not isinstance(companies_raw, dict):
                continue
            for company_name, company_data in companies_raw.items():
                if not isinstance(company_data, dict):
                    continue
                normalized = self._normalize_company_name(company_name)
                if normalized not in merged_companies:
                    merged_companies[normalized] = {'assets': []}

                for asset in company_data.get('assets', []):
                    existing = self._find_matching_asset(
                        merged_companies[normalized]['assets'], asset)
                    if existing:
                        # Merge into existing asset
                        for update in asset.get('key_updates', []):
                            if update and update not in existing.get('key_updates', []):
                                existing.setdefault('key_updates', []).append(update)
                        for data in asset.get('clinical_data', []):
                            if data and data not in existing.get('clinical_data', []):
                                existing.setdefault('clinical_data', []).append(data)
                        if asset.get('stage') and not existing.get('stage'):
                            existing['stage'] = asset['stage']
                        if asset.get('target') and not existing.get('target'):
                            existing['target'] = asset['target']
                        if asset.get('payload') and not existing.get('payload'):
                            existing['payload'] = asset['payload']
                        for ind in asset.get('indications', []):
                            if ind and ind not in existing.get('indications', []):
                                existing.setdefault('indications', []).append(ind)
                    else:
                        merged_companies[normalized]['assets'].append(asset)

            # Merge deals
            for deal in result.get('deals', []):
                if deal and not self._is_duplicate_deal(merged_deals, deal):
                    merged_deals.append(deal)

            # Merge regulatory events
            for event in result.get('regulatory_events', []):
                if event and not self._is_duplicate_event(merged_regulatory, event):
                    merged_regulatory.append(event)

        # Build summary lists
        companies_mentioned = sorted(merged_companies.keys())
        assets_mentioned = sorted(set(
            asset.get('name', '')
            for company_data in merged_companies.values()
            for asset in company_data.get('assets', [])
            if asset.get('name')
        ))

        return {
            'companies': merged_companies,
            'deals': merged_deals,
            'regulatory_events': merged_regulatory,
            'companies_mentioned': companies_mentioned,
            'assets_mentioned': assets_mentioned,
        }

    COMPANY_ALIASES = {
        'az': 'AstraZeneca',
        'astra zeneca': 'AstraZeneca',
        'astrazeneca': 'AstraZeneca',
        'daiichi-sankyo': 'Daiichi Sankyo',
        'daiichisankyo': 'Daiichi Sankyo',
        'daiichi sankyo': 'Daiichi Sankyo',
        'bms': 'Bristol-Myers Squibb',
        'bristol myers squibb': 'Bristol-Myers Squibb',
        'bristol-myers squibb': 'Bristol-Myers Squibb',
        'roche/genentech': 'Roche',
        'genentech': 'Roche',
        'msd': 'Merck & Co.',
        'merck & co': 'Merck & Co.',
        'merck & co.': 'Merck & Co.',
        'lilly': 'Eli Lilly',
        'eli lilly': 'Eli Lilly',
        'j&j': 'Johnson & Johnson',
        'janssen': 'Johnson & Johnson',
        'pfizer inc': 'Pfizer',
        'pfizer inc.': 'Pfizer',
        'gsk': 'GlaxoSmithKline',
        'glaxo smith kline': 'GlaxoSmithKline',
        'glaxosmithkline': 'GlaxoSmithKline',
        'sanofi-aventis': 'Sanofi',
        'abbvie': 'AbbVie',
        'seagen': 'Pfizer',
    }

    def _normalize_company_name(self, name: str) -> str:
        """Normalize company name for deduplication (case-insensitive)."""
        name = name.strip()
        return self.COMPANY_ALIASES.get(name.lower(), name)

    def _find_matching_asset(self, existing_assets: List[Dict], new_asset: Dict) -> Optional[Dict]:
        """Find if an asset already exists in the list (by name matching)."""
        new_name = (new_asset.get('name') or '').lower().strip()
        if not new_name:
            return None
        for asset in existing_assets:
            existing_name = (asset.get('name') or '').lower().strip()
            if existing_name == new_name:
                return asset
            # Only allow substring matching for code-like names (e.g., DS-8201, T-DXd)
            shorter = new_name if len(new_name) <= len(existing_name) else existing_name
            longer = existing_name if len(new_name) <= len(existing_name) else new_name
            if (len(shorter) >= 5 and shorter in longer
                    and self.ASSET_CODE_RE.match(shorter)):
                return asset
        return None

    def _is_duplicate_deal(self, existing_deals: List[Dict], new_deal: Dict) -> bool:
        new_parties = frozenset(new_deal.get('parties', []))
        new_details = (new_deal.get('details') or '').lower()[:80]
        for deal in existing_deals:
            if frozenset(deal.get('parties', [])) == new_parties:
                existing_details = (deal.get('details') or '').lower()
                if (new_details and
                        (new_details in existing_details or existing_details[:80] in new_details)):
                    return True
        return False

    def _is_duplicate_event(self, existing_events: List[Dict], new_event: Dict) -> bool:
        new_key = (
            (new_event.get('company') or '').lower(),
            (new_event.get('asset') or '').lower(),
            (new_event.get('event') or '').lower()[:50],
        )
        for event in existing_events:
            existing_key = (
                (event.get('company') or '').lower(),
                (event.get('asset') or '').lower(),
                (event.get('event') or '').lower()[:50],
            )
            if new_key == existing_key:
                return True
        return False

    # ─── Executive Summary ───────────────────────────────────────────────────

    def generate_executive_summary(self, ci_data: Dict, doc_type: str, title: str) -> str:
        """Generate a CI-focused executive summary from structured extraction data."""
        if not self._track_api_call():
            return ""

        client = self._get_client()

        # Build compact structured payload (much smaller than 8K raw text)
        summary_input = {
            'companies': ci_data.get('companies_mentioned', []),
            'assets': ci_data.get('assets_mentioned', []),
            'deals': ci_data.get('deals', []),
            'regulatory': ci_data.get('regulatory_events', []),
            'top_updates': [],
        }
        # Include the most important updates (first 2 per company, max 20 total)
        for company_data in list(ci_data.get('companies', {}).values())[:15]:
            for asset in company_data.get('assets', [])[:2]:
                for update in asset.get('key_updates', [])[:1]:
                    summary_input['top_updates'].append(
                        f"{asset.get('name', '?')}: {update}")
                for data in asset.get('clinical_data', [])[:1]:
                    summary_input['top_updates'].append(
                        f"{asset.get('name', '?')} data: {data}")
        summary_input['top_updates'] = summary_input['top_updates'][:20]

        payload = json.dumps(summary_input, separators=(',', ':'))

        prompt = f"""Summarize this ADC competitive intelligence for "{title}" ({doc_type}) in 3-5 bullet points.
Focus on: most significant developments, clinical data, regulatory decisions, deals.
Be specific and data-driven (include numbers, company names, drug names).

{payload}"""

        try:
            response = self._api_call_with_retry(
                lambda: client.messages.create(
                    model=self.VISION_MODEL,
                    max_tokens=1024,
                    messages=[{"role": "user", "content": prompt}]
                ),
                description="Executive summary"
            )
            return response.content[0].text.strip()
        except Exception as e:
            logger.error(f"  Summary generation failed: {e}")
            return ""

    # ─── Quality Scoring ─────────────────────────────────────────────────────

    def calculate_quality_score(self, full_text: str, ci_data: Dict,
                                has_summary: bool, metadata: Dict) -> Dict:
        """Calculate quality score for CI document processing."""
        scores = {}

        # Text extraction (0-10)
        text_len = len(full_text)
        if text_len > 20000:
            scores['text_extraction'] = 10.0
        elif text_len > 10000:
            scores['text_extraction'] = 8.0
        elif text_len > 5000:
            scores['text_extraction'] = 6.0
        elif text_len > 1000:
            scores['text_extraction'] = 4.0
        else:
            scores['text_extraction'] = 2.0

        # Entity extraction (0-10) — most important for CI
        companies = ci_data.get('companies', {})
        total_assets = sum(len(c.get('assets', [])) for c in companies.values())
        if total_assets >= 20:
            scores['entity_extraction'] = 10.0
        elif total_assets >= 10:
            scores['entity_extraction'] = 8.0
        elif total_assets >= 5:
            scores['entity_extraction'] = 6.0
        elif total_assets >= 1:
            scores['entity_extraction'] = 4.0
        else:
            scores['entity_extraction'] = 2.0

        # Metadata completeness (0-10)
        meta_score = 0.0
        if metadata.get('title') and metadata['title'] != 'Untitled':
            meta_score += 2.5
        if metadata.get('date'):
            meta_score += 2.5
        if metadata.get('source') and metadata['source'] != 'Unknown':
            meta_score += 2.5
        if metadata.get('doc_type') and metadata['doc_type'] != 'other':
            meta_score += 2.5
        scores['metadata_completeness'] = meta_score

        # Structure quality (0-10)
        deals = ci_data.get('deals', [])
        regulatory = ci_data.get('regulatory_events', [])
        struct_score = 5.0
        if deals:
            struct_score += 2.0
        if regulatory:
            struct_score += 2.0
        if total_assets > 0:
            assets_with_stage = sum(
                1 for c in companies.values()
                for a in c.get('assets', []) if a.get('stage')
            )
            completeness = assets_with_stage / max(total_assets, 1)
            struct_score += completeness * 1.0
        scores['structure_quality'] = min(10.0, struct_score)

        # Summary quality (0-10)
        scores['summary_quality'] = 8.0 if has_summary else 3.0

        weights = {
            'text_extraction': 0.25,
            'entity_extraction': 0.30,
            'metadata_completeness': 0.15,
            'structure_quality': 0.15,
            'summary_quality': 0.15,
        }
        overall = sum(scores[k] * weights[k] for k in weights)
        overall = round(min(10.0, max(0.0, overall)), 1)

        if overall >= 8.0:
            assessment = "Excellent"
        elif overall >= 5.5:
            assessment = "Good"
        elif overall >= 4.0:
            assessment = "Fair"
        else:
            assessment = "Poor"

        return {
            'overall': overall,
            'assessment': assessment,
            'components': scores,
        }

    # ─── Output Generation ───────────────────────────────────────────────────

    def _make_title_slug(self, title: str, max_words: int = 6) -> str:
        words = re.findall(r'[A-Za-z0-9]+', title)
        stop_words = {'a', 'an', 'the', 'of', 'in', 'on', 'at', 'to', 'for',
                      'and', 'or', 'but', 'is', 'are', 'was', 'were', 'with', 'by'}
        meaningful = [w for w in words if w.lower() not in stop_words]
        if not meaningful:
            meaningful = words
        slug = '_'.join(meaningful[:max_words]).lower()
        return slug if slug else 'untitled'

    def generate_output_filename(self, doc_type: str, source: str,
                                  title: str, date: str) -> str:
        title_slug = self._make_title_slug(title)
        source_slug = re.sub(r'[^a-zA-Z0-9]', '', source).lower()[:15]
        date_slug = date or 'undated'

        template = self.NAMING_SCHEMES.get(self.naming, self.NAMING_SCHEMES['default'])
        filename = template.format(
            source_type=doc_type,
            title_slug=title_slug,
            date=date_slug,
            source=source_slug,
        )
        return sanitize_filename(f"{filename}.md")

    def generate_markdown(self, metadata: Dict, ci_data: Dict, summary: str,
                          quality: Dict, content: Dict) -> str:
        """Generate the final CI-structured markdown document."""
        lines = []

        # YAML frontmatter
        lines.append('---')
        lines.append(f'title: "{self._escape_yaml(metadata["title"])}"')
        lines.append(f'filename: "{self._escape_yaml(metadata["filename"])}"')
        lines.append(f'source_type: {metadata["doc_type"]}')
        if metadata.get('date'):
            lines.append(f'date: "{metadata["date"]}"')
        lines.append(f'source: "{self._escape_yaml(metadata["source"])}"')
        lines.append('classification: confidential')

        companies_mentioned = ci_data.get('companies_mentioned', [])
        if companies_mentioned:
            lines.append('companies_mentioned:')
            for c in companies_mentioned[:50]:
                lines.append(f'  - "{self._escape_yaml(c)}"')

        assets_mentioned = ci_data.get('assets_mentioned', [])
        if assets_mentioned:
            lines.append('assets_mentioned:')
            for a in assets_mentioned[:50]:
                lines.append(f'  - "{self._escape_yaml(a)}"')

        lines.append(f'extraction_method: {content["extraction_method"]}')
        lines.append(f'vision_model: {self.VISION_MODEL}')
        lines.append(f'processing_date: {datetime.now().strftime("%Y-%m-%d")}')
        lines.append(f'quality_overall: {quality["overall"]}')
        lines.append(f'quality_assessment: {quality["assessment"]}')
        lines.append(f'total_pages: {content["total_pages"]}')
        lines.append(f'api_calls_used: {self._api_call_count}')
        lines.append('---')
        lines.append('')

        # Title
        lines.append(f'# {metadata["title"]}')
        lines.append('')
        header_parts = []
        if metadata.get('date'):
            header_parts.append(f'**Date**: {metadata["date"]}')
        header_parts.append(f'**Source**: {metadata["source"]}')
        header_parts.append(f'**Type**: {metadata["doc_type"]}')
        lines.append('  |  '.join(header_parts))
        lines.append('')
        lines.append('---')
        lines.append('')

        # Executive Summary
        if summary:
            lines.append('## Executive Summary')
            lines.append('')
            lines.append(summary)
            lines.append('')
            lines.append('---')
            lines.append('')

        # Companies & Assets
        companies = ci_data.get('companies', {})
        if companies:
            lines.append('## Companies & Assets')
            lines.append('')
            for company_name in sorted(companies.keys()):
                company = companies[company_name]
                assets = company.get('assets', [])
                lines.append(f'### {company_name}')
                lines.append('')

                if not assets:
                    lines.append('*Mentioned in document (no specific asset details extracted)*')
                    lines.append('')
                    continue

                for asset in assets:
                    # Asset header with metadata
                    header_parts = []
                    if asset.get('target'):
                        header_parts.append(f'Target: {asset["target"]}')
                    if asset.get('payload'):
                        header_parts.append(f'Payload: {asset["payload"]}')
                    suffix = f' ({", ".join(header_parts)})' if header_parts else ''
                    asset_name = asset.get('name', 'Unknown asset')
                    lines.append(f'#### {asset_name}{suffix}')
                    lines.append('')

                    if asset.get('modality'):
                        lines.append(f'- **Modality**: {asset["modality"]}')
                    if asset.get('stage'):
                        lines.append(f'- **Stage**: {asset["stage"]}')
                    if asset.get('indications'):
                        lines.append(f'- **Indications**: {", ".join(asset["indications"])}')
                    for update in asset.get('key_updates', []):
                        lines.append(f'- **Update**: {update}')
                    for data in asset.get('clinical_data', []):
                        lines.append(f'- **Clinical Data**: {data}')
                    lines.append('')

            lines.append('---')
            lines.append('')

        # Deals & Partnerships
        deals = ci_data.get('deals', [])
        if deals:
            lines.append('## Key Deals & Partnerships')
            lines.append('')
            for deal in deals:
                parties = ', '.join(deal.get('parties', []))
                deal_type = deal.get('type', 'deal').title()
                details = deal.get('details', '')
                lines.append(f'- **{deal_type}** ({parties}): {details}')
            lines.append('')
            lines.append('---')
            lines.append('')

        # Regulatory Updates
        reg_events = ci_data.get('regulatory_events', [])
        if reg_events:
            lines.append('## Regulatory Updates')
            lines.append('')
            for event in reg_events:
                company = event.get('company', '?')
                asset = event.get('asset', '?')
                evt = event.get('event', '')
                lines.append(f'- **{company}** / {asset}: {evt}')
            lines.append('')
            lines.append('---')
            lines.append('')

        # Quality Assessment
        lines.append('## Quality Assessment')
        lines.append('')
        lines.append(f'**Overall**: {quality["overall"]}/10 — {quality["assessment"]}')
        lines.append('')
        for comp, score in quality.get('components', {}).items():
            label = comp.replace('_', ' ').title()
            lines.append(f'- {label}: {score:.1f}/10')
        lines.append('')
        lines.append('---')
        lines.append('')

        # Processing Metadata
        lines.append('## Processing Metadata')
        lines.append('')
        lines.append(f'- **Extraction Method**: {content["extraction_method"]}')
        lines.append(f'- **Total Pages**: {content["total_pages"]}')
        lines.append(f'- **Text Length**: {content["total_chars"]:,} characters')
        lines.append(f'- **Companies Found**: {len(ci_data.get("companies", {}))}')
        lines.append(f'- **Assets Found**: {len(ci_data.get("assets_mentioned", []))}')
        lines.append(f'- **API Calls Used**: {self._api_call_count}')
        lines.append(f'- **Processing Date**: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
        lines.append('')

        return '\n'.join(lines)

    # ─── Single Document Processing ──────────────────────────────────────────

    def process_single_document(self, filepath: Path, skip_existing: bool = True) -> bool:
        """Process a single CI document to structured markdown."""
        self._api_call_count = 0

        logger.info(f"\n{'=' * 60}")
        logger.info(f"Processing: {filepath.name}")
        logger.info(f"{'=' * 60}")

        # File size check
        file_size_mb = filepath.stat().st_size / (1024 * 1024)
        if file_size_mb > self.MAX_FILE_SIZE_MB:
            logger.error(f"  File too large: {file_size_mb:.1f}MB — skipping")
            return False

        # Step 1: Document type detection
        doc_type, source = self.detect_document_type(filepath)
        date = self._extract_date_from_filename(filepath)
        title = self._extract_title_from_filename(filepath, doc_type)
        logger.info(f"  Type: {doc_type} | Source: {source} | Date: {date or 'unknown'}")

        # Check skip-existing
        output_filename = self.generate_output_filename(doc_type, source, title, date)
        if skip_existing and output_filename.replace('.md', '') in self._existing_files:
            logger.info(f"  Already processed — skipping")
            return True

        # Step 2: Content extraction
        logger.info("  [1/4] Extracting content...")
        if filepath.suffix.lower() == '.pptx':
            content = self.extract_pptx_content(filepath)
        else:
            content = self.extract_pdf_content(filepath)

        if not content['full_text'].strip():
            logger.error(f"  Failed to extract text content")
            self._log_processing(filepath.name, doc_type, source, 0, 0, 0, 0, 0.0, "FAILED")
            return False

        logger.info(f"  Pages: {content['total_pages']} | Chars: {content['total_chars']:,}")

        # Step 3: Company & asset extraction
        logger.info("  [2/4] Extracting companies and assets...")
        ci_data = self.extract_companies_and_assets(
            content['full_text'], doc_type, source)

        companies_count = len(ci_data.get('companies', {}))
        assets_count = len(ci_data.get('assets_mentioned', []))
        deals_count = len(ci_data.get('deals', []))
        logger.info(f"  Companies: {companies_count} | Assets: {assets_count} | Deals: {deals_count}")

        # Step 4: Executive summary
        logger.info("  [3/4] Generating executive summary...")
        summary = self.generate_executive_summary(ci_data, doc_type, title)

        # Step 5: Quality scoring
        logger.info("  [4/4] Quality assessment...")
        metadata = {
            'title': title,
            'filename': filepath.name,
            'doc_type': doc_type,
            'source': source,
            'date': date,
        }
        quality = self.calculate_quality_score(
            content['full_text'], ci_data, bool(summary), metadata)
        logger.info(f"  Quality: {quality['overall']}/10 — {quality['assessment']}")

        # Generate and save markdown
        markdown = self.generate_markdown(metadata, ci_data, summary, quality, content)

        output_path = self.output_dir / output_filename
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown)
        self._existing_files.add(output_path.stem)

        self._log_processing(
            filepath.name, doc_type, source, content['total_pages'],
            content['total_chars'], companies_count, assets_count,
            quality['overall'], "SAVED"
        )

        logger.info(f"  SAVED: {output_filename}")
        logger.info(f"  API calls used: {self._api_call_count}")
        return True

    # ─── Batch Processing ────────────────────────────────────────────────────

    def process_all_documents(self, skip_existing: bool = True):
        """Process all CI documents in the input folder."""
        logger.info("=" * 70)
        logger.info("COMPETITIVE INTELLIGENCE PROCESSING PIPELINE")
        logger.info("=" * 70)
        logger.info(f"Input folder: {self.input_folder}")
        logger.info(f"Output directory: {self.output_dir}")
        logger.info(f"Recursive: {self.recursive}")
        logger.info(f"Budget: {self._budget_cap or 'unlimited'} calls per document")
        logger.info(f"Parallel workers: {self.MAX_CONCURRENT_CHUNKS}")
        logger.info("=" * 70)

        files = self.list_ci_files()
        if not files:
            logger.warning("No documents found")
            return

        success_count = 0
        fail_count = 0

        for i, filepath in enumerate(files, 1):
            logger.info(f"\n[{i}/{len(files)}] {filepath.name}")
            try:
                if self.process_single_document(filepath, skip_existing=skip_existing):
                    success_count += 1
                else:
                    fail_count += 1
            except Exception as e:
                logger.error(f"  EXCEPTION: {e}")
                fail_count += 1

        logger.info(f"\n{'=' * 70}")
        logger.info(f"PIPELINE COMPLETE: {success_count} succeeded, {fail_count} failed "
                    f"out of {len(files)}")
        logger.info(f"Processing log: {self.processing_log_path}")
        logger.info(f"{'=' * 70}")


def main():
    parser = argparse.ArgumentParser(
        description='Process competitive intelligence documents to structured markdown',
        epilog='Supports PDF and PPTX. Extracts companies, assets, and key events.'
    )
    parser.add_argument('--input', type=str, required=True,
                        help='Path to folder containing CI documents')
    parser.add_argument('--output', type=str, default='output_ci',
                        help='Output directory for markdown files (default: output_ci)')
    parser.add_argument('--single', type=str, default=None,
                        help='Process a single file only (provide full path)')
    parser.add_argument('--recursive', action='store_true',
                        help='Recursively search subfolders')
    parser.add_argument('--no-skip', action='store_true',
                        help='Reprocess documents even if output exists')
    parser.add_argument('--naming', default='default',
                        choices=['default', 'dated', 'source'],
                        help="Output filename scheme: "
                             "'default' = ci_{type}_{title}; "
                             "'dated' = ci_{date}_{type}_{title}; "
                             "'source' = ci_{source}_{date}_{title}")
    parser.add_argument('--budget', type=int, default=30,
                        help='Max API calls per document (default: 30, 0=unlimited)')
    parser.add_argument('--verbose', action='store_true',
                        help='Enable debug logging')

    args = parser.parse_args()

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    try:
        validate_path(args.input, must_exist=True, allow_dir=True, allow_file=False)
        validate_output_path(args.output)
        if args.single:
            validate_path(args.single, must_exist=True, allow_file=True, allow_dir=False)
    except (ValueError, FileNotFoundError) as e:
        logger.error(f"Path validation failed: {e}")
        return

    pipeline = CIPipeline(
        input_folder=args.input,
        output_dir=args.output,
        recursive=args.recursive,
        naming=args.naming,
        budget=args.budget,
    )

    if args.single:
        single_path = Path(args.single)
        pipeline.process_single_document(single_path, skip_existing=not args.no_skip)
    else:
        pipeline.process_all_documents(skip_existing=not args.no_skip)


if __name__ == '__main__':
    main()
